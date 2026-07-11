import os
import re
import io
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import streamlit as st
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# ReportLab imports for PDF generation
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from sklearn.metrics import (
    accuracy_score,
    precision_recall_fscore_support,
    confusion_matrix,
    ConfusionMatrixDisplay,
    cohen_kappa_score
)

from src.config import OUTPUT_FILE, YOUTUBE_VIDEO_URL, MAX_COMMENTS, HISTORY_DIR, APP_MODE
from src.downloader import fetch_youtube_comments, get_video_title, extract_video_id, get_video_context
from src.lexicon_analyzer import LexiconSentimentAnalyzer
from src.llm_analyzer import LLMSentimentAnalyzer

# Conditional GSheets Connection for Production Mode
conn = None
if APP_MODE == "production":
    try:
        from streamlit_gsheets import GSheetsConnection
        conn = st.connection("gsheets", type=GSheetsConnection)
    except Exception as e:
        # Display warning if the spreadsheet fails to load
        st.warning(f"Gagal menghubungkan ke Google Sheets: {e}")


def interpret_kappa(score):
    if score < 0:
        return "Tidak ada kesepakatan"
    elif score <= 0.20:
        return "Sangat Lemah"
    elif score <= 0.40:
        return "Lemah"
    elif score <= 0.60:
        return "Sedang"
    elif score <= 0.80:
        return "Kuat"
    else:
        return "Sangat Kuat"


def upgrade_dataframe_schema(df):
    """
    Upgrades older CSV schemas to the new dual-mode schema containing columns:
    'LLM Sentiment Global', 'LLM Reason Global', 'LLM Sentiment Video', 'LLM Reason Video'
    """
    df = df.copy()
    
    # Ensure new columns exist
    if "LLM Sentiment Global" not in df.columns:
        df["LLM Sentiment Global"] = None
    if "LLM Reason Global" not in df.columns:
        df["LLM Reason Global"] = ""
    if "LLM Sentiment Video" not in df.columns:
        df["LLM Sentiment Video"] = None
    if "LLM Reason Video" not in df.columns:
        df["LLM Reason Video"] = ""
        
    # If old columns are present, map them based on 'Analysis Mode'
    if "LLM Sentiment" in df.columns:
        for idx, row in df.iterrows():
            mode = str(row.get("Analysis Mode", "")).strip().lower()
            val = row.get("LLM Sentiment", "")
            reason = row.get("LLM Reason", "")
            
            # Map values
            if "video" in mode:
                if pd.isna(df.at[idx, "LLM Sentiment Video"]) or str(df.at[idx, "LLM Sentiment Video"]).strip() in ["", "nan", "None"]:
                    df.at[idx, "LLM Sentiment Video"] = val
                    df.at[idx, "LLM Reason Video"] = reason
            else:
                if pd.isna(df.at[idx, "LLM Sentiment Global"]) or str(df.at[idx, "LLM Sentiment Global"]).strip() in ["", "nan", "None"]:
                    df.at[idx, "LLM Sentiment Global"] = val
                    df.at[idx, "LLM Reason Global"] = reason
                    
    # Fill remaining NaNs with fallback value 'netral'
    df["LLM Sentiment Global"] = df["LLM Sentiment Global"].fillna("netral").astype(str)
    df["LLM Sentiment Video"] = df["LLM Sentiment Video"].fillna("netral").astype(str)
    df["LLM Reason Global"] = df["LLM Reason Global"].fillna("").astype(str)
    df["LLM Reason Video"] = df["LLM Reason Video"].fillna("").astype(str)

    # Ensure Lexicon columns exist
    if "Lexicon Sentiment" not in df.columns:
        df["Lexicon Sentiment"] = "netral"
    if "Lexicon Score" not in df.columns:
        df["Lexicon Score"] = 0
    df["Lexicon Sentiment"] = df["Lexicon Sentiment"].fillna("netral").astype(str)
        
    # Ensure Ground Truth exists
    if "Ground Truth" not in df.columns:
        df["Ground Truth"] = ""
    df["Ground Truth"] = df["Ground Truth"].fillna("")
        
    return df


    return df


def render_evaluation_metrics(df_eval, llm_col_sentiment, llm_col_reason, mode_title):
    # Filter rows with Ground Truth
    df_eval_filtered = df_eval.dropna(subset=["Ground Truth"]).copy()
    df_eval_filtered = df_eval_filtered[df_eval_filtered["Ground Truth"].astype(str).str.strip().str.lower().isin(["positif", "negatif", "netral"])]
    
    if len(df_eval_filtered) == 0:
        st.warning(f"Belum ada Ground Truth yang diisi untuk {mode_title}. Silakan isi beberapa baris pada kolom Ground Truth di tabel atas untuk memunculkan evaluasi metrik akurasi.", icon=":material/warning:")
        return None
        
    total_gt = len(df_eval_filtered)
    st.success(
        f"Menghitung performa {mode_title} berdasarkan **{total_gt}** komentar yang telah dilabeli Ground Truth (Evaluasi seluruh kelas: Positif, Negatif, dan Netral).",
        icon=":material/check_circle:"
    )
    
    y_true = df_eval_filtered["Ground Truth"].str.strip().str.lower()
    y_lexicon = df_eval_filtered["Lexicon Sentiment"].str.strip().str.lower()
    y_llm = df_eval_filtered[llm_col_sentiment].str.strip().str.lower()
    
    # Standard ML Scores
    lex_acc = accuracy_score(y_true, y_lexicon)
    lex_prec, lex_rec, lex_f1, _ = precision_recall_fscore_support(y_true, y_lexicon, average='macro', zero_division=0)
    kappa_lex = cohen_kappa_score(y_true, y_lexicon)
    
    llm_acc = accuracy_score(y_true, y_llm)
    llm_prec, llm_rec, llm_f1, _ = precision_recall_fscore_support(y_true, y_llm, average='macro', zero_division=0)
    kappa_llm = cohen_kappa_score(y_true, y_llm)
    
    # Calculate SEMANTIKA Points System (all 3 classes)
    lex_points = 0
    llm_points = 0
    lex_correct = 0
    llm_correct = 0
    
    # Breakdown counters
    lex_correct_pos = 0
    lex_correct_neg = 0
    lex_correct_net = 0
    
    llm_correct_pos = 0
    llm_correct_neg = 0
    llm_correct_net = 0
    
    total_pos = 0
    total_neg = 0
    total_net = 0
    
    for idx, row in df_eval_filtered.iterrows():
        gt = str(row["Ground Truth"]).strip().lower()
        lex = str(row["Lexicon Sentiment"]).strip().lower()
        llm = str(row[llm_col_sentiment]).strip().lower()
        
        if gt == "positif":
            total_pos += 1
        elif gt == "negatif":
            total_neg += 1
        elif gt == "netral":
            total_net += 1
            
        # Lexicon
        if lex == gt:
            lex_points += 1
            lex_correct += 1
            if gt == "positif":
                lex_correct_pos += 1
            elif gt == "negatif":
                lex_correct_neg += 1
            elif gt == "netral":
                lex_correct_net += 1
        else:
            lex_points -= 1
            
        # LLM
        if llm == gt:
            llm_points += 1
            llm_correct += 1
            if gt == "positif":
                llm_correct_pos += 1
            elif gt == "negatif":
                llm_correct_neg += 1
            elif gt == "netral":
                llm_correct_net += 1
        else:
            llm_points -= 1

    total_eval_comments = len(df_eval_filtered)

    # Render Points Comparison UI Cards
    col_pts1, col_pts2 = st.columns(2)
    with col_pts1:
        st.markdown(
            f"""
            <div class="point-card lexicon-card">
                <h3>POIN PERFORMA LEXICON</h3>
                <div style="font-size: 3rem; font-weight: 800; margin: 5px 0;">{lex_points}</div>
                <div style="font-size: 1rem; font-weight: 600; opacity: 0.9; margin-bottom: 5px;">Total Benar: {lex_correct} / {total_eval_comments} Komentar</div>
                <div style="font-size: 0.85rem; opacity: 0.8; line-height: 1.4; margin-bottom: 8px; text-align: left; padding-left: 20%;">
                    • Positif: {lex_correct_pos} / {total_pos}<br/>
                    • Negatif: {lex_correct_neg} / {total_neg}<br/>
                    • Netral: {lex_correct_net} / {total_net}
                </div>
                <p style="margin-top: 5px; font-size: 0.8rem;">Metode Sastrawi + InSet Lexicon</p>
            </div>
            """,
            unsafe_allow_html=True
        )
        
    with col_pts2:
        st.markdown(
            f"""
            <div class="point-card llm-card">
                <h3>POIN PERFORMA LLM ({mode_title})</h3>
                <div style="font-size: 3rem; font-weight: 800; margin: 5px 0;">{llm_points}</div>
                <div style="font-size: 1rem; font-weight: 600; opacity: 0.9; margin-bottom: 5px;">Total Benar: {llm_correct} / {total_eval_comments} Komentar</div>
                <div style="font-size: 0.85rem; opacity: 0.8; line-height: 1.4; margin-bottom: 8px; text-align: left; padding-left: 20%;">
                    • Positif: {llm_correct_pos} / {total_pos}<br/>
                    • Negatif: {llm_correct_neg} / {total_neg}<br/>
                    • Netral: {llm_correct_net} / {total_net}
                </div>
                <p style="margin-top: 5px; font-size: 0.8rem;">NVIDIA NIM ({st.session_state.llm_model.split('/')[-1]})</p>
            </div>
            """,
            unsafe_allow_html=True
        )
        
    # Render Points Donut Charts for Accuracy
    def pct_val_fmt(pct, allvals):
        absolute = int(round(pct/100.*np.sum(allvals)))
        return f"{pct:.1f}%\n({absolute} data)"
        
    fig_perf, (ax_p1, ax_p2) = plt.subplots(1, 2, figsize=(8, 3))
    perf_colors = ["#2ecc71", "#e74c3c"]
    
    # Lexicon Donut
    lex_incorrect = total_eval_comments - lex_correct
    if total_eval_comments > 0:
        sizes = [lex_correct, lex_incorrect]
        labels = ["Benar", "Salah"]
        ax_p1.pie(sizes, labels=labels, autopct=lambda pct: pct_val_fmt(pct, sizes), startangle=90, colors=[perf_colors[0] if l == "Benar" else perf_colors[1] for l in labels], pctdistance=0.70, textprops=dict(color="black", weight="bold", fontsize=8))
        ax_p1.set_title("Akurasi Lexicon (Benar vs Salah)", fontsize=9, fontweight="bold")
        centre_circle = plt.Circle((0,0), 0.50, fc='white')
        ax_p1.add_artist(centre_circle)
        
    # LLM Donut
    llm_incorrect = total_eval_comments - llm_correct
    if total_eval_comments > 0:
        sizes = [llm_correct, llm_incorrect]
        labels = ["Benar", "Salah"]
        ax_p2.pie(sizes, labels=labels, autopct=lambda pct: pct_val_fmt(pct, sizes), startangle=90, colors=[perf_colors[0] if l == "Benar" else perf_colors[1] for l in labels], pctdistance=0.70, textprops=dict(color="black", weight="bold", fontsize=8))
        ax_p2.set_title(f"Akurasi LLM {mode_title}", fontsize=9, fontweight="bold")
        centre_circle = plt.Circle((0,0), 0.50, fc='white')
        ax_p2.add_artist(centre_circle)
        
    plt.tight_layout()
    st.pyplot(fig_perf)
    plt.close()
    
    # Point Rules Explanation
    st.markdown("""
    <div class="info-box">
        <strong>Sistem Poin Komparasi SEMANTIKA:</strong><br/>
        Sistem poin di atas dihitung dengan ketentuan:
        <ul>
            <li>Setiap data Ground Truth yang bernilai <strong>positif</strong>, <strong>negatif</strong>, atau <strong>netral</strong> akan dievaluasi.</li>
            <li>Jika tebakan model <strong>Benar (sesuai Ground Truth)</strong> &rarr; Model mendapatkan <strong>+1 Poin</strong>.</li>
            <li>Jika tebakan model <strong>Salah</strong> &rarr; Model dikurangi <strong>-1 Poin</strong>.</li>
        </ul>
    </div>
    """, unsafe_allow_html=True)
    
    # Metrics Table
    st.markdown("##### Tabel Rincian Metrik Klasifikasi")
    metrics_data = {
        "Metrik": ["Akurasi", "Presisi (Macro)", "Recall (Macro)", "F1-Score (Macro)", "Cohen's Kappa (Kesepakatan)"],
        "Lexicon-Based": [f"{lex_acc*100:.1f}%", f"{lex_prec*100:.1f}%", f"{lex_rec*100:.1f}%", f"{lex_f1*100:.1f}%", f"{kappa_lex:.3f} ({interpret_kappa(kappa_lex)})"],
        "LLM-Based": [f"{llm_acc*100:.1f}%", f"{llm_prec*100:.1f}%", f"{llm_rec*100:.1f}%", f"{llm_f1*100:.1f}%", f"{kappa_llm:.3f} ({interpret_kappa(kappa_llm)})"]
    }
    st.table(pd.DataFrame(metrics_data))
    
    return {
        "lex_acc": lex_acc, "lex_prec": lex_prec, "lex_rec": lex_rec, "lex_f1": lex_f1, "kappa_lex": kappa_lex, "lex_points": lex_points,
        "llm_acc": llm_acc, "llm_prec": llm_prec, "llm_rec": llm_rec, "llm_f1": llm_f1, "kappa_llm": kappa_llm, "llm_points": llm_points,
        "total_pos": total_pos, "total_neg": total_neg, "total_net": total_net
    }


def render_mode_visualizations(mode_key, llm_col_sentiment, llm_col_reason, mode_title):
    df_eval = st.session_state.df.dropna(subset=["Ground Truth"]).copy()
    df_eval = df_eval[df_eval["Ground Truth"].astype(str).str.strip().str.lower().isin(["positif", "negatif", "netral"])]
    
    if len(df_eval) == 0:
        st.warning("Belum ada data Ground Truth terisi. Silakan isi Ground Truth pada tabel untuk memunculkan visualisasi.", icon=":material/warning:")
        return

    y_true = df_eval["Ground Truth"].str.strip().str.lower()
    y_lexicon = df_eval["Lexicon Sentiment"].str.strip().str.lower()
    y_llm = df_eval[llm_col_sentiment].str.strip().str.lower()
    
    # Calculate scores on the fly for metrics display
    lex_acc = accuracy_score(y_true, y_lexicon)
    lex_prec, lex_rec, lex_f1, _ = precision_recall_fscore_support(y_true, y_lexicon, average='macro', zero_division=0)
    kappa_lex = cohen_kappa_score(y_true, y_lexicon)
    
    llm_acc = accuracy_score(y_true, y_llm)
    llm_prec, llm_rec, llm_f1, _ = precision_recall_fscore_support(y_true, y_llm, average='macro', zero_division=0)
    kappa_llm = cohen_kappa_score(y_true, y_llm)

    # Tab Layout for Visualizations
    tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
        ":material/pie_chart: Sebaran Sentimen (Donut Charts)", 
        ":material/equalizer: Performa Klasifikasi (Metrics)",
        ":material/bar_chart: Perbandingan Metrik (Bar Chart)",
        ":material/timeline: Tren Sentimen (Timeline)",
        ":material/short_text: Frekuensi Kata (Top Words)",
        ":material/category: Pemodelan Topik (Topic Modeling)"
    ])
    
    # Tab 1: Donut Charts
    with tab1:
        st.markdown(f"### Perbandingan Sebaran Sentimen ({mode_title})")
        sentiment_labels = ["positif", "negatif", "netral"]
        color_map = {"positif": "#2ecc71", "negatif": "#e74c3c", "netral": "#95a5a6"}
        
        def get_sizes_and_colors(series):
            counts = series.value_counts()
            sizes = []
            colors = []
            labels = []
            for label in sentiment_labels:
                count = counts.get(label, 0)
                if count > 0:
                    sizes.append(count)
                    colors.append(color_map[label])
                    labels.append(label.capitalize())
            return sizes, colors, labels
            
        gt_sizes, gt_colors, gt_labels = get_sizes_and_colors(y_true)
        lex_sizes, lex_colors, lex_labels = get_sizes_and_colors(y_lexicon)
        llm_sizes, llm_colors, llm_labels = get_sizes_and_colors(y_llm)
        
        fig_donut, (ax_d1, ax_d2, ax_d3) = plt.subplots(1, 3, figsize=(18, 6))
        
        # Donut 1: Ground Truth
        if gt_sizes:
            ax_d1.pie(gt_sizes, labels=gt_labels, autopct='%1.1f%%', startangle=90, colors=gt_colors, pctdistance=0.75, textprops=dict(color="black", weight="bold"))
            centre_circle = plt.Circle((0,0), 0.50, fc='white')
            ax_d1.add_artist(centre_circle)
            ax_d1.set_title("Sebaran Ground Truth", fontsize=12, weight="bold")
        else:
            ax_d1.text(0.5, 0.5, 'Tidak ada data', ha='center', va='center')
            
        # Donut 2: Lexicon
        if lex_sizes:
            ax_d2.pie(lex_sizes, labels=lex_labels, autopct='%1.1f%%', startangle=90, colors=lex_colors, pctdistance=0.75, textprops=dict(color="black", weight="bold"))
            centre_circle = plt.Circle((0,0), 0.50, fc='white')
            ax_d2.add_artist(centre_circle)
            ax_d2.set_title("Sebaran Lexicon-based", fontsize=12, weight="bold")
        else:
            ax_d2.text(0.5, 0.5, 'Tidak ada data', ha='center', va='center')
            
        # Donut 3: LLM
        if llm_sizes:
            ax_d3.pie(llm_sizes, labels=llm_labels, autopct='%1.1f%%', startangle=90, colors=llm_colors, pctdistance=0.75, textprops=dict(color="black", weight="bold"))
            centre_circle = plt.Circle((0,0), 0.50, fc='white')
            ax_d3.add_artist(centre_circle)
            ax_d3.set_title(f"Sebaran LLM-based\n({st.session_state.llm_model.split('/')[-1]})", fontsize=12, weight="bold")
        else:
            ax_d3.text(0.5, 0.5, 'Tidak ada data', ha='center', va='center')
            
        plt.tight_layout()
        st.pyplot(fig_donut)
        plt.close()

    # Tab 2: Performa Klasifikasi & Akurasi
    with tab2:
        st.markdown(f"### Performa Klasifikasi & Akurasi ({mode_title})")
        col_m1, col_m2 = st.columns(2)
        with col_m1:
            st.markdown("#### :material/book: Performa Lexicon (Sastrawi + InSet)")
            st.write(f"- **Accuracy (Akurasi)**: {lex_acc * 100:.2f}%")
            st.write(f"- **Precision (Presisi)**: {lex_prec * 100:.2f}%")
            st.write(f"- **Recall (Sensitivitas)**: {lex_rec * 100:.2f}%")
            st.write(f"- **F1-Score**: {lex_f1 * 100:.2f}%")
            st.write(f"- **Cohen's Kappa**: `{kappa_lex:.4f}` ({interpret_kappa(kappa_lex)})")
            
        with col_m2:
            st.markdown(f"#### :material/psychology: Performa LLM ({st.session_state.llm_model.split('/')[-1]})")
            st.write(f"- **Accuracy (Akurasi)**: {llm_acc * 100:.2f}%")
            st.write(f"- **Precision (Presisi)**: {llm_prec * 100:.2f}%")
            st.write(f"- **Recall (Sensitivitas)**: {llm_rec * 100:.2f}%")
            st.write(f"- **F1-Score**: {llm_f1 * 100:.2f}%")
            st.write(f"- **Cohen's Kappa**: `{kappa_llm:.4f}` ({interpret_kappa(kappa_llm)})")

    # Tab 3: Metrics Comparison Bar Chart
    with tab3:
        st.markdown(f"### Perbandingan Metrik Evaluasi (Lexicon vs LLM - {mode_title})")
        metrics = ['Akurasi', 'Presisi', 'Sensitivitas (Recall)', 'F1-Score']
        lex_scores = [lex_acc * 100, lex_prec * 100, lex_rec * 100, lex_f1 * 100]
        llm_scores = [llm_acc * 100, llm_prec * 100, llm_rec * 100, llm_f1 * 100]
        
        x = np.arange(len(metrics))
        width = 0.35
        
        fig_metrics, ax_m = plt.subplots(figsize=(10, 5))
        rects1 = ax_m.bar(x - width/2, lex_scores, width, label='Lexicon-based', color='#3498db')
        rects2 = ax_m.bar(x + width/2, llm_scores, width, label=f'LLM-based ({st.session_state.llm_model.split("/")[-1]})', color='#2ecc71')
        
        ax_m.set_ylabel('Skor (%)', weight="bold")
        ax_m.set_title(f'Perbandingan Metrik {mode_title}', weight="bold", fontsize=12)
        ax_m.set_xticks(x)
        ax_m.set_xticklabels(metrics, weight="bold")
        ax_m.set_ylim(0, 110)
        ax_m.legend()
        
        def autolabel(rects):
            for rect in rects:
                height = rect.get_height()
                ax_m.annotate(f'{height:.1f}%',
                            xy=(rect.get_x() + rect.get_width() / 2, height),
                            xytext=(0, 3),
                            textcoords="offset points",
                            ha='center', va='bottom', weight="bold", fontsize=9)
                            
        autolabel(rects1)
        autolabel(rects2)
        plt.tight_layout()
        st.pyplot(fig_metrics)
        plt.close()
        
        st.markdown("---")
        col_k1, col_k2 = st.columns(2)
        with col_k1:
            st.info(f"**Cohen's Kappa Lexicon vs GT:**\n\n`{kappa_lex:.4f}` $\\rightarrow$ Kesepakatan **{interpret_kappa(kappa_lex)}**", icon=":material/book:")
        with col_k2:
            st.info(f"**Cohen's Kappa LLM vs GT:**\n\n`{kappa_llm:.4f}` $\\rightarrow$ Kesepakatan **{interpret_kappa(kappa_llm)}**", icon=":material/psychology:")

    # Tab 4: Tren Sentimen terhadap Waktu Rilis Komentar
    with tab4:
        st.markdown("### Tren Sentimen terhadap Waktu Rilis Komentar")
        if "Timestamp" in df_eval.columns and df_eval["Timestamp"].notna().any():
            try:
                df_time = df_eval.copy()
                df_time["Datetime"] = pd.to_datetime(df_time["Timestamp"], unit='s', errors='coerce')
                df_time = df_time.dropna(subset=["Datetime"])
                
                if len(df_time) > 0:
                    df_time["Date"] = df_time["Datetime"].dt.date
                    time_pivot = df_time.groupby(["Date", "Ground Truth"]).size().unstack(fill_value=0)
                    for col in ["positif", "negatif", "netral"]:
                        if col not in time_pivot.columns:
                            time_pivot[col] = 0
                    time_pivot = time_pivot.sort_index()
                    
                    fig_time, ax_t = plt.subplots(figsize=(12, 5))
                    if "positif" in time_pivot.columns:
                        ax_t.plot(time_pivot.index, time_pivot["positif"], marker='o', color='#2ecc71', label='Positif', linewidth=2.5)
                    if "negatif" in time_pivot.columns:
                        ax_t.plot(time_pivot.index, time_pivot["negatif"], marker='x', color='#e74c3c', label='Negatif', linewidth=2.5)
                    if "netral" in time_pivot.columns:
                        ax_t.plot(time_pivot.index, time_pivot["netral"], marker='s', color='#95a5a6', label='Netral', linewidth=2)
                    
                    ax_t.set_ylabel('Jumlah Komentar', weight="bold")
                    ax_t.set_xlabel('Tanggal', weight="bold")
                    ax_t.set_title('Tren Perkembangan Sentimen Komentar (Ground Truth)', weight="bold", fontsize=12)
                    ax_t.grid(True, linestyle='--', alpha=0.5)
                    ax_t.legend()
                    plt.xticks(rotation=45)
                    plt.tight_layout()
                    st.pyplot(fig_time)
                    plt.close()
                else:
                    st.info("Data timestamp tidak valid.")
            except Exception as e:
                st.error(f"Gagal memproses tren waktu: {e}")
        else:
            st.info("Komentar tidak memiliki data timestamp yang valid.")

    # Tab 5: Analisis Frekuensi Kata Terpopuler
    with tab5:
        st.markdown("### Analisis Frekuensi Kata Kunci Terpopuler")
        def get_top_words(df_subset, top_n=10):
            word_counts = {}
            for text in df_subset["Cleaned Comment"].dropna().astype(str):
                for word in text.split():
                    if len(word) > 1:
                        word_counts[word] = word_counts.get(word, 0) + 1
            sorted_words = sorted(word_counts.items(), key=lambda x: x[1], reverse=True)
            return sorted_words[:top_n]
            
        df_pos = df_eval[df_eval["Ground Truth"] == "positif"]
        df_neg = df_eval[df_eval["Ground Truth"] == "negatif"]
        top_pos = get_top_words(df_pos)
        top_neg = get_top_words(df_neg)
        
        col_pos, col_neg = st.columns(2)
        with col_pos:
            st.markdown("#### Kata Kunci pada Komentar Positif")
            if top_pos:
                words, counts = zip(*top_pos)
                fig_p, ax_p = plt.subplots(figsize=(6, 4))
                ax_p.barh(words[::-1], counts[::-1], color='#2ecc71')
                ax_p.set_title("Top Kata Komentar Positif", weight="bold")
                plt.tight_layout()
                st.pyplot(fig_p)
                plt.close()
            else:
                st.info("Belum ada data kata kunci positif.")
        with col_neg:
            st.markdown("#### Kata Kunci pada Komentar Negatif")
            if top_neg:
                words, counts = zip(*top_neg)
                fig_n, ax_n = plt.subplots(figsize=(6, 4))
                ax_n.barh(words[::-1], counts[::-1], color='#e74c3c')
                ax_n.set_title("Top Kata Komentar Negatif", weight="bold")
                plt.tight_layout()
                st.pyplot(fig_n)
                plt.close()
            else:
                st.info("Belum ada data kata kunci negatif.")

    # Tab 6: Pemodelan Topik (TF-IDF + KMeans)
    with tab6:
        st.markdown("### Pemodelan Topik (Topic Modeling)")
        clean_comments = df_eval["Cleaned Comment"].dropna().astype(str).tolist()
        original_comments = df_eval["Original Comment"].dropna().astype(str).tolist()
        
        if len(clean_comments) >= 5:
            try:
                from sklearn.feature_extraction.text import TfidfVectorizer
                from sklearn.cluster import KMeans
                vectorizer = TfidfVectorizer(max_features=500, min_df=1, stop_words=None)
                X = vectorizer.fit_transform(clean_comments)
                n_clusters = min(3, len(clean_comments))
                kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init='auto')
                kmeans.fit(X)
                order_centroids = kmeans.cluster_centers_.argsort()[:, ::-1]
                terms = vectorizer.get_feature_names_out()
                
                st.markdown(f"Berhasil mendeteksi **{n_clusters}** kelompok topik utama pembicaraan:")
                for cluster_idx in range(n_clusters):
                    top_words = [terms[ind] for ind in order_centroids[cluster_idx, :5]]
                    topic_keywords = ", ".join(top_words)
                    cluster_comments = [original_comments[j] for j, label in enumerate(kmeans.labels_) if label == cluster_idx]
                    sample_comments = cluster_comments[:3]
                    st.markdown(f"**Topik {cluster_idx + 1}:** `{topic_keywords}` ({len(cluster_comments)} komentar)")
                    for sc in sample_comments:
                        st.caption(f"- \"{sc[:120]}...\"")
                    st.markdown(" ")
            except Exception as e:
                st.error(f"Gagal melakukan clustering topik: {e}")
        else:
            st.info("Dibutuhkan minimal 5 komentar Ground Truth untuk mendeteksi kelompok topik secara akurat.")


def render_mode_tab_content(mode_key, llm_col_sentiment, llm_col_reason, mode_title):
    # Retrieve base data
    df_temp = st.session_state.df.copy()
    
    # Check if LLM results are missing for this mode
    missing_results = df_temp[llm_col_sentiment].isna() | (df_temp[llm_col_sentiment].astype(str).str.strip() == "")
    if missing_results.sum() > 0:
        st.warning(f"Hasil analisis untuk mode **{mode_title}** belum lengkap ({missing_results.sum()}/{len(df_temp)} komentar kosong). Silakan jalankan analisis ulang untuk melengkapinya.")
        
    # We display the metadata context in Video Context mode
    if mode_key == "video":
        if "video_context" not in st.session_state or st.session_state.video_context is None:
            if st.session_state.video_url:
                try:
                    with st.spinner("Mengambil transkrip/konteks video dari YouTube..."):
                        st.session_state.video_context = get_video_context(st.session_state.video_url)
                except Exception:
                    st.session_state.video_context = "Gagal mengambil transkrip/konteks video."
            else:
                st.session_state.video_context = "Teks transkrip/konteks video tidak tersedia."
        with st.expander("🎥 Lihat Metadata & Transkrip Subtitel Video (Konteks)", expanded=False):
            st.markdown("Informasi berikut diekstraksi secara otomatis dari YouTube untuk mencocokkan emosi dan relevansi komentar:")
            if st.session_state.video_context:
                st.text_area("Konteks Video (Title, Description, Transcript):", value=st.session_state.video_context, height=250, disabled=True, key=f"video_ctx_area_{mode_key}")
            else:
                st.info("Konteks video tidak tersedia.")

    # Sort & Filters for this mode's table
    agree_mask = df_temp["Lexicon Sentiment"] == df_temp[llm_col_sentiment]
    total_comments = len(df_temp)
    agreed_comments = agree_mask.sum()
    pct_agreement = (agreed_comments / total_comments * 100) if total_comments > 0 else 0
    
    empty_gt_mask = df_temp["Ground Truth"].isna() | (df_temp["Ground Truth"].astype(str).str.strip() == "")
    fillable_count = empty_gt_mask.sum()
    
    col_filter, col_actions = st.columns([3, 2])
    with col_filter:
        show_mismatch = st.checkbox(f" Hanya tampilkan data Mismatch (Lexicon vs LLM berbeda)", value=False, key=f"filter_mismatch_{mode_key}")
        st.markdown(f"Tingkat kesepakatan model (Lexicon & LLM sama): **{pct_agreement:.1f}%** ({agreed_comments}/{total_comments} komentar)")
        sort_by = st.selectbox(
            " Urutkan Komentar",
            options=["Bawaan (YouTube)", "Likes Terbanyak", "Waktu Terbaru"],
            key=f"comment_sort_select_{mode_key}",
            help="Urutkan komentar pada tabel."
        )
    with col_actions:
        with st.expander("🛠 Fitur Pengisian Cepat (Quick Labeling)", expanded=False):
            st.write(f"Komentar yang bisa diisi otomatis (kosong): **{fillable_count}**")
            btn_col1, btn_col2 = st.columns(2)
            with btn_col1:
                if st.button("Isi Otomatis", use_container_width=True, key=f"btn_fill_{mode_key}", help="Mengisi Ground Truth kosong dengan hasil prediksi DeepSeek V4 Pro"):
                    if fillable_count > 0:
                        with st.status("Mengisi Ground Truth menggunakan DeepSeek V4 Pro...", expanded=True) as status:
                            if st.session_state.llm_model == "deepseek-ai/deepseek-v4-pro":
                                st.session_state.df.loc[empty_gt_mask, "Ground Truth"] = st.session_state.df.loc[empty_gt_mask, llm_col_sentiment]
                            else:
                                # Fetch on the fly
                                rows_to_fill = st.session_state.df[empty_gt_mask]
                                comments_to_analyze = []
                                for idx, row in rows_to_fill.iterrows():
                                    comments_to_analyze.append({
                                        "comment_id": row["Comment ID"],
                                        "text": row["Original Comment"]
                                    })
                                
                                ds_analyzer = LLMSentimentAnalyzer(model="deepseek-ai/deepseek-v4-pro")
                                ds_results = []
                                batch_size = 20
                                num_batches = (len(comments_to_analyze) - 1) // batch_size + 1
                                for batch_idx in range(0, len(comments_to_analyze), batch_size):
                                    batch = comments_to_analyze[batch_idx:batch_idx+batch_size]
                                    status.write(f"   - Memproses Batch {batch_idx//batch_size + 1}/{num_batches}...")
                                    video_context = get_video_context(st.session_state.video_url) if mode_key == "video" else None
                                    batch_results = ds_analyzer.analyze_batch(batch, video_context=video_context)
                                    ds_results.extend(batch_results)
                                
                                for r in ds_results:
                                    cid = r["comment_id"]
                                    sentiment = r["llm_sentiment"]
                                    st.session_state.df.loc[st.session_state.df["Comment ID"] == cid, "Ground Truth"] = sentiment
                                    
                            st.session_state.df.to_csv(OUTPUT_FILE, index=False, encoding="utf-8-sig")
                            video_id = extract_video_id(st.session_state.video_url)
                            if video_id:
                                safe_title = make_safe_filename(st.session_state.video_title)
                                history_filename = f"[{video_id}] {safe_title}.csv"
                                history_path = os.path.join(HISTORY_DIR, history_filename)
                                st.session_state.df.to_csv(history_path, index=False, encoding="utf-8-sig")
                                if APP_MODE == "production":
                                    sync_video_to_gsheets(video_id, st.session_state.video_title, st.session_state.video_url, st.session_state.df)
                            st.rerun()
            with btn_col2:
                if st.button("Kosongkan GT", use_container_width=True, key=f"btn_clear_{mode_key}", help="Menghapus semua label Ground Truth"):
                    st.session_state.df["Ground Truth"] = ""
                    st.session_state.df.to_csv(OUTPUT_FILE, index=False, encoding="utf-8-sig")
                    video_id = extract_video_id(st.session_state.video_url)
                    if video_id:
                        safe_title = make_safe_filename(st.session_state.video_title)
                        history_filename = f"[{video_id}] {safe_title}.csv"
                        history_path = os.path.join(HISTORY_DIR, history_filename)
                        st.session_state.df.to_csv(history_path, index=False, encoding="utf-8-sig")
                        if APP_MODE == "production":
                            sync_video_to_gsheets(video_id, st.session_state.video_title, st.session_state.video_url, st.session_state.df)
                    st.success("Ground Truth dikosongkan!")
                    st.rerun()

    # Prepare Display DataFrame
    display_df = df_temp.copy()
    
    # Sort
    if sort_by == "Likes Terbanyak" and "Likes" in display_df.columns:
        display_df["Likes"] = pd.to_numeric(display_df["Likes"], errors='coerce').fillna(0).astype(int)
        display_df = display_df.sort_values(by="Likes", ascending=False)
    elif sort_by == "Waktu Terbaru" and "Timestamp" in display_df.columns:
        display_df["Timestamp"] = pd.to_numeric(display_df["Timestamp"], errors='coerce').fillna(0)
        display_df = display_df.sort_values(by="Timestamp", ascending=False)
        
    if show_mismatch:
        display_df = display_df[display_df["Lexicon Sentiment"] != display_df[llm_col_sentiment]]

    # Map table columns
    table_display_df = display_df.copy()
    table_display_df["LLM Sentiment"] = table_display_df[llm_col_sentiment]
    table_display_df["LLM Reason"] = table_display_df[llm_col_reason]

    tab_view_tbl, tab_edit_tbl = st.tabs([":material/visibility: Tampilan Tabel Berwarna", ":material/edit: Edit Ground Truth"])
    
    with tab_view_tbl:
        def style_table_row(row):
            styles = pd.Series("", index=row.index)
            def get_color(val):
                val_lower = str(val).strip().lower()
                if val_lower == "positif":
                    return "background-color: #C6EFCE; color: #006100; font-weight: bold;"
                elif val_lower == "negatif":
                    return "background-color: #FFC7CE; color: #9C0006; font-weight: bold;"
                elif val_lower == "netral":
                    return "background-color: #E2E3E5; color: #383D41;"
                return ""
            
            lex = row.get("Lexicon Sentiment")
            llm = row.get("LLM Sentiment")
            gt = row.get("Ground Truth")
            if "Lexicon Sentiment" in row.index:
                styles["Lexicon Sentiment"] = get_color(lex)
            if "LLM Sentiment" in row.index:
                styles["LLM Sentiment"] = get_color(llm)
            if "Ground Truth" in row.index:
                if str(lex).strip().lower() != str(llm).strip().lower():
                    styles["Ground Truth"] = "background-color: #FFE699; color: #7F6000; font-weight: bold;"
                else:
                    styles["Ground Truth"] = get_color(gt)
            return styles

        styled_df = table_display_df.style.apply(style_table_row, axis=1)
        st.dataframe(
            styled_df,
            column_config={
                "No": st.column_config.NumberColumn("No", width="small"),
                "Author": st.column_config.TextColumn("Penulis", width="medium"),
                "Original Comment": st.column_config.TextColumn("Komentar Asli", width="large"),
                "Cleaned Comment": st.column_config.TextColumn("Komentar Bersih (Stemmed)", width="medium"),
                "Likes": st.column_config.NumberColumn("Likes", width="small"),
                "Time Description": st.column_config.TextColumn("Waktu", width="small"),
                "Lexicon Sentiment": st.column_config.TextColumn("Sentimen Lexicon", width="small"),
                "LLM Sentiment": st.column_config.TextColumn("Sentimen LLM", width="small"),
                "LLM Reason": st.column_config.TextColumn("Alasan LLM", width="medium"),
                "Ground Truth": st.column_config.TextColumn("Ground Truth", width="small"),
            },
            column_order=["No", "Author", "Original Comment", "Cleaned Comment", "Likes", "Time Description", "Lexicon Sentiment", "LLM Sentiment", "LLM Reason", "Ground Truth"],
            use_container_width=True,
            hide_index=True,
            key=f"df_view_{mode_key}"
        )
        
    with tab_edit_tbl:
        st.info("Gunakan tabel di bawah ini untuk menentukan sentimen sebenarnya (Ground Truth) lewat dropdown pilihan.", icon=":material/info:")
        edited_df = st.data_editor(
            table_display_df,
            column_config={
                "Ground Truth": st.column_config.SelectboxColumn(
                    "Ground Truth",
                    options=["positif", "negatif", "netral", ""],
                    required=False
                ),
                "No": st.column_config.NumberColumn("No", width="small", disabled=True),
                "Author": st.column_config.TextColumn("Penulis", width="medium", disabled=True),
                "Original Comment": st.column_config.TextColumn("Komentar Asli", width="large", disabled=True),
                "Cleaned Comment": st.column_config.TextColumn("Komentar Bersih (Stemmed)", width="medium", disabled=True),
                "Likes": st.column_config.NumberColumn("Likes", width="small", disabled=True),
                "Time Description": st.column_config.TextColumn("Waktu", width="small", disabled=True),
                "Lexicon Sentiment": st.column_config.TextColumn("Sentimen Lexicon", width="small", disabled=True),
                "LLM Sentiment": st.column_config.TextColumn("Sentimen LLM", width="small", disabled=True),
                "LLM Reason": st.column_config.TextColumn("Alasan LLM", width="medium", disabled=True),
            },
            column_order=["No", "Author", "Original Comment", "Cleaned Comment", "Likes", "Time Description", "Lexicon Sentiment", "LLM Sentiment", "LLM Reason", "Ground Truth"],
            use_container_width=True,
            key=f"editor_{mode_key}",
            num_rows="fixed"
        )
        
        if not edited_df.equals(table_display_df):
            st.session_state.df.loc[edited_df.index, "Ground Truth"] = edited_df["Ground Truth"]
            st.session_state.df.to_csv(OUTPUT_FILE, index=False, encoding="utf-8-sig")
            video_id = extract_video_id(st.session_state.video_url)
            if video_id:
                safe_title = make_safe_filename(st.session_state.video_title)
                history_filename = f"[{video_id}] {safe_title}.csv"
                history_path = os.path.join(HISTORY_DIR, history_filename)
                st.session_state.df.to_csv(history_path, index=False, encoding="utf-8-sig")
                if APP_MODE == "production":
                    sync_video_to_gsheets(video_id, st.session_state.video_title, st.session_state.video_url, st.session_state.df)
            st.rerun()

    # Download Button Section
    st.markdown(" ")
    col_dl1, col_dl2, col_dl3 = st.columns(3)
    with col_dl1:
        # Prepare df for export
        df_export = st.session_state.df.copy()
        df_export["LLM Sentiment"] = df_export[llm_col_sentiment]
        df_export["LLM Reason"] = df_export[llm_col_reason]
        df_export["Analysis Mode"] = mode_title
        
        excel_data = convert_df_to_excel(df_export, st.session_state.video_title, st.session_state.video_url)
        st.download_button(
            label=":material/download: Ekspor Excel (.xlsx)",
            data=excel_data,
            file_name=f"semantika_hasil_{extract_video_id(st.session_state.video_url)}_{mode_key}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
            key=f"dl_excel_{mode_key}"
        )
    with col_dl2:
        df_export = st.session_state.df.copy()
        df_export["LLM Sentiment"] = df_export[llm_col_sentiment]
        df_export["LLM Reason"] = df_export[llm_col_reason]
        df_export["Analysis Mode"] = mode_title
        
        pdf_data = convert_df_to_pdf(df_export, st.session_state.video_title, st.session_state.video_url)
        st.download_button(
            label=":material/download: Ekspor PDF (.pdf)",
            data=pdf_data,
            file_name=f"semantika_hasil_{extract_video_id(st.session_state.video_url)}_{mode_key}.pdf",
            mime="application/pdf",
            use_container_width=True,
            key=f"dl_pdf_{mode_key}"
        )
    with col_dl3:
        df_export = st.session_state.df.copy()
        df_export["LLM Sentiment"] = df_export[llm_col_sentiment]
        df_export["LLM Reason"] = df_export[llm_col_reason]
        df_export["Analysis Mode"] = mode_title
        
        pptx_data = convert_df_to_pptx(
            df_export,
            st.session_state.video_title,
            st.session_state.video_url,
            mode_title,
            st.session_state.llm_model
        )
        st.download_button(
            label=":material/download: Ekspor Presentasi (.pptx)",
            data=pptx_data,
            file_name=f"semantika_presentasi_{extract_video_id(st.session_state.video_url)}_{mode_key}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            key=f"dl_pptx_{mode_key}"
        )

    # Section 3: Live Evaluation
    st.markdown("---")
    st.subheader(f":material/trending_up: Evaluasi Performa Real-Time ({mode_title})")
    
    df_eval = st.session_state.df.copy()
    eval_stats = render_evaluation_metrics(df_eval, llm_col_sentiment, llm_col_reason, mode_title)

    # Section 4: Execution Benchmark
    st.markdown("---")
    st.subheader(":material/timer: Perbandingan Kecepatan Eksekusi (Benchmark)")
    lex_t = st.session_state.get("lexicon_time")
    llm_t = st.session_state.get(f"llm_time_{mode_key}")
    
    if lex_t is not None and llm_t is not None:
        col_bench1, col_bench2 = st.columns(2)
        with col_bench1:
            lex_avg = (lex_t / total_comments) * 1000 if total_comments > 0 else 0
            st.metric(
                label="Total Durasi Analisis Lexicon (Offline/Lokal)",
                value=f"{lex_t:.4f} detik",
                delta=f"{lex_avg:.2f} ms / komentar",
                delta_color="normal"
            )
        with col_bench2:
            llm_avg = (llm_t / total_comments) * 1000 if total_comments > 0 else 0
            st.metric(
                label=f"Total Durasi Analisis LLM ({mode_title})",
                value=f"{llm_t:.2f} detik",
                delta=f"{llm_avg:.0f} ms / komentar",
                delta_color="inverse"
            )
        if lex_t > 0:
            speedup = llm_t / lex_t
            st.info(f":material/bolt: **Hasil Benchmark Kecepatan:** Metode Lexicon berjalan **{speedup:.1f}x lebih cepat** dibandingkan metode LLM karena diproses secara lokal tanpa latency jaringan.", icon=":material/info:")
    else:
        st.info("Informasi waktu eksekusi benchmark hanya tersedia untuk video yang baru dianalisis pada sesi aktif saat ini.", icon=":material/info:")

    # Section 5: Visualizations specific to this mode
    st.markdown("---")
    st.subheader(f":material/bar_chart: Tab Analitik & Visualisasi ({mode_title})")
    
    render_mode_visualizations(mode_key, llm_col_sentiment, llm_col_reason, mode_title)


# Set Streamlit Page Config
st.set_page_config(
    page_title="SEMANTIKA - YouTube Sentiment Analysis Dashboard",
    page_icon=":material/analytics:",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for Premium Design Aesthetics
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=Outfit:wght@400;500;600;700;800&display=swap');

    /* Global Fonts & Body background layering */
    .stApp {
        background-color: #f8fafc;
        font-family: 'Inter', sans-serif;
    }
    
    /* Title styling */
    h1, h2, h3 {
        font-family: 'Outfit', sans-serif !important;
        color: #0f172a !important;
        font-weight: 700 !important;
        letter-spacing: -0.5px !important;
    }
    
    /* Premium Metric Card */
    .metric-card {
        background: #ffffff;
        border: 1px solid #edf2f7;
        border-radius: 16px;
        padding: 24px;
        box-shadow: 0 4px 20px -2px rgba(50, 50, 93, 0.04), 0 2px 8px -1px rgba(0, 0, 0, 0.02);
        transition: all 0.25s cubic-bezier(0.4, 0, 0.2, 1);
        margin-bottom: 20px;
    }
    .metric-card:hover {
        transform: translateY(-3px);
        box-shadow: 0 12px 24px -10px rgba(50, 50, 93, 0.08), 0 4px 12px -5px rgba(0, 0, 0, 0.03);
        border-color: #e2e8f0;
    }
    
    .metric-title {
        font-size: 0.8rem;
        color: #94a3b8;
        font-weight: 600;
        text-transform: uppercase;
        letter-spacing: 0.8px;
        margin-bottom: 8px;
    }
    
    .metric-value {
        font-size: 2.25rem;
        font-weight: 700;
        color: #0f172a;
        font-family: 'Outfit', sans-serif;
    }

    /* Point board cards with sleek gradients */
    .point-card {
        border-radius: 16px;
        padding: 28px;
        text-align: center;
        color: white;
        box-shadow: 0 10px 25px -5px rgba(0, 0, 0, 0.05);
        transition: transform 0.2s ease;
    }
    .point-card:hover {
        transform: scale(1.02);
    }
    .lexicon-card {
        background: linear-gradient(135deg, #3b82f6, #1d4ed8);
    }
    .llm-card {
        background: linear-gradient(135deg, #10b981, #047857);
    }
    
    /* Info box styling */
    .info-box {
        background-color: #ffffff;
        border-left: 4px solid #3b82f6;
        padding: 16px 20px;
        border-radius: 4px 16px 16px 4px;
        box-shadow: 0 2px 10px rgba(0,0,0,0.02);
        margin-bottom: 20px;
        border-top: 1px solid #f1f5f9;
        border-right: 1px solid #f1f5f9;
        border-bottom: 1px solid #f1f5f9;
    }

    /* Streamlit Tab Custom Styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        background-color: #f1f5f9;
        padding: 6px;
        border-radius: 12px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 40px;
        white-space: pre-wrap;
        background-color: transparent;
        border-radius: 8px;
        color: #64748b;
        font-weight: 500;
        font-size: 0.9rem;
        border: none;
        padding: 0 16px;
    }
    .stTabs [aria-selected="true"] {
        background-color: #ffffff !important;
        color: #0f172a !important;
        font-weight: 600 !important;
        box-shadow: 0 2px 8px -1px rgba(0,0,0,0.05);
    }

    /* Sidebar background & border */
    section[data-testid="stSidebar"] {
        background-color: #ffffff;
        border-right: 1px solid #e2e8f0;
    }

    /* Buttons premium styling */
    .stButton>button {
        border-radius: 10px !important;
        font-weight: 500 !important;
        transition: all 0.2s ease !important;
    }
</style>
""", unsafe_allow_html=True)

# Initialize Session State
if "df" not in st.session_state:
    st.session_state.df = None
if "video_title" not in st.session_state:
    st.session_state.video_title = ""
if "video_url" not in st.session_state:
    st.session_state.video_url = ""
if "llm_model" not in st.session_state:
    st.session_state.llm_model = "meta/llama-3.1-8b-instruct"
if "detected_lang" not in st.session_state:
    st.session_state.detected_lang = "id"
if "youtube_url_widget" not in st.session_state:
    st.session_state.youtube_url_widget = YOUTUBE_VIDEO_URL
if "loaded_history_file" not in st.session_state:
    st.session_state.loaded_history_file = ""
if "lexicon_time" not in st.session_state:
    st.session_state.lexicon_time = None
if "llm_time" not in st.session_state:
    st.session_state.llm_time = None
if "analysis_mode" not in st.session_state:
    st.session_state.analysis_mode = "Konteks Global"
if "video_context" not in st.session_state:
    st.session_state.video_context = None

# Auto-load existing results if CSV exists
if st.session_state.df is None and os.path.exists(OUTPUT_FILE):
    try:
        df_loaded = pd.read_csv(OUTPUT_FILE)
        df_loaded = upgrade_dataframe_schema(df_loaded)
        required_cols = ["No", "Comment ID", "Author", "Original Comment", "Cleaned Comment", "Lexicon Sentiment", "Lexicon Score", "Ground Truth"]
        if all(col in df_loaded.columns for col in required_cols):
            df_loaded["Ground Truth"] = df_loaded["Ground Truth"].fillna("")
            if "LLM Model" in df_loaded.columns:
                st.session_state.llm_model = str(df_loaded["LLM Model"].iloc[0])
            else:
                df_loaded["LLM Model"] = "meta/llama-3.1-8b-instruct"
                st.session_state.llm_model = "meta/llama-3.1-8b-instruct"
            if "Language" in df_loaded.columns:
                st.session_state.detected_lang = str(df_loaded["Language"].iloc[0])
            else:
                df_loaded["Language"] = "id"
                st.session_state.detected_lang = "id"
            
            st.session_state.analysis_mode = "Konteks Global"
                
            if "Lexicon Time" in df_loaded.columns:
                st.session_state.lexicon_time = float(df_loaded["Lexicon Time"].dropna().iloc[0]) if len(df_loaded["Lexicon Time"].dropna()) > 0 else len(df_loaded) * 0.0015
            else:
                st.session_state.lexicon_time = len(df_loaded) * 0.0015
                
            if "LLM Time" in df_loaded.columns:
                st.session_state.llm_time = float(df_loaded["LLM Time"].dropna().iloc[0]) if len(df_loaded["LLM Time"].dropna()) > 0 else len(df_loaded) * 0.15
            else:
                st.session_state.llm_time = len(df_loaded) * 0.15
                
            st.session_state.df = df_loaded
            st.session_state.video_url = YOUTUBE_VIDEO_URL
            st.session_state.video_title = get_video_title(YOUTUBE_VIDEO_URL)
            st.session_state.youtube_url_widget = YOUTUBE_VIDEO_URL
    except Exception:
        pass

# Helper to convert DataFrame to a beautifully styled Excel file in memory
def convert_df_to_excel(df, video_title, video_url):
    output = io.BytesIO()
    
    # Rename columns to match local labels
    df_export = df.copy()
    df_export = df_export.rename(columns={
        "No": "No",
        "Author": "Penulis",
        "Original Comment": "Komentar Asli",
        "Cleaned Comment": "Komentar Bersih (Stemmed)",
        "Lexicon Sentiment": "Sentimen Lexicon",
        "LLM Sentiment": "Sentimen LLM",
        "LLM Reason": "Alasan LLM",
        "Ground Truth": "Ground Truth"
    })
    
    # Keep only target columns in the exported spreadsheet
    cols_to_keep = ["No", "Penulis", "Komentar Asli", "Komentar Bersih (Stemmed)", "Sentimen Lexicon", "Sentimen LLM", "Alasan LLM", "Ground Truth"]
    df_export = df_export[[col for col in cols_to_keep if col in df_export.columns]]
    
    # Column Widths mapping (defined early for row height estimates)
    col_widths = {
        "No": 6,
        "Penulis": 18,
        "Komentar Asli": 50,
        "Komentar Bersih (Stemmed)": 40,
        "Sentimen Lexicon": 18,
        "Sentimen LLM": 18,
        "Alasan LLM": 40,
        "Ground Truth": 18
    }
    
    # Write to Excel in memory using openpyxl, start data at Row 7 (startrow=6)
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df_export.to_excel(writer, sheet_name='Analisis Sentimen', index=False, startrow=6)
        workbook = writer.book
        worksheet = writer.sheets['Analisis Sentimen']
        
        # Write Report Title Headers at the top (Rows 1-5)
        worksheet.cell(row=1, column=1, value="SEMANTIKA - Laporan Analisis Sentimen Komentar YouTube")
        worksheet.cell(row=1, column=1).font = Font(name='Arial', size=15, bold=True, color='1F4E79')
        
        worksheet.cell(row=2, column=1, value=f"Judul Video: {video_title}")
        worksheet.cell(row=2, column=1).font = Font(name='Arial', size=10, bold=True)
        
        worksheet.cell(row=3, column=1, value=f"Link Video: {video_url}")
        worksheet.cell(row=3, column=1).font = Font(name='Arial', size=10, color='2563EB', underline='single')
        
        worksheet.cell(row=4, column=1, value=f"Model LLM: {st.session_state.llm_model}")
        worksheet.cell(row=4, column=1).font = Font(name='Arial', size=10, bold=True)
        
        lang_label = "Inggris (EN)" if st.session_state.detected_lang == "en" else "Indonesia (ID)"
        worksheet.cell(row=5, column=1, value=f"Bahasa Terdeteksi: {lang_label}")
        worksheet.cell(row=5, column=1).font = Font(name='Arial', size=10, bold=True)
        
        # Color palettes & Font settings for Table
        header_fill = PatternFill(start_color='1F4E79', end_color='1F4E79', fill_type='solid') # Slate Blue
        header_font = Font(name='Arial', size=11, bold=True, color='FFFFFF')
        
        zebra_fill = PatternFill(start_color='F2F4F8', end_color='F2F4F8', fill_type='solid') # Alternating light gray/blue
        white_fill = PatternFill(start_color='FFFFFF', end_color='FFFFFF', fill_type='solid')
        
        # Soft fills and colors for sentiment states
        positif_fill = PatternFill(start_color='C6EFCE', end_color='C6EFCE', fill_type='solid') # Light Green
        positif_font = Font(name='Arial', size=10, color='006100', bold=True)
        
        negatif_fill = PatternFill(start_color='FFC7CE', end_color='FFC7CE', fill_type='solid') # Light Red
        negatif_font = Font(name='Arial', size=10, color='9C0006', bold=True)
        
        netral_fill = PatternFill(start_color='E2E3E5', end_color='E2E3E5', fill_type='solid')  # Light Gray
        netral_font = Font(name='Arial', size=10, color='383D41', bold=False)
        
        mismatch_fill = PatternFill(start_color='FFE699', end_color='FFE699', fill_type='solid') # Warning Yellow
        mismatch_font = Font(name='Arial', size=10, color='7F6000', bold=True)
        
        thin_border = Border(
            left=Side(style='thin', color='D3D3D3'),
            right=Side(style='thin', color='D3D3D3'),
            top=Side(style='thin', color='D3D3D3'),
            bottom=Side(style='thin', color='D3D3D3')
        )
        
        align_center = Alignment(horizontal='center', vertical='center')
        align_left = Alignment(horizontal='left', vertical='top', wrap_text=True)
        
        # Style Table Header (Row 7)
        for col_idx in range(1, len(df_export.columns) + 1):
            cell = worksheet.cell(row=7, column=col_idx)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = align_center
            cell.border = thin_border
            
        # Style Table Data Rows (Row 8 onwards) & Calculate Row Heights dynamically
        for row_idx in range(8, worksheet.max_row + 1):
            # Apply zebra striping
            row_fill = zebra_fill if row_idx % 2 == 0 else white_fill
            
            max_lines = 1
            for col_idx, col_name in enumerate(df_export.columns, start=1):
                cell = worksheet.cell(row=row_idx, column=col_idx)
                cell.fill = row_fill
                cell.border = thin_border
                
                # Column specific alignments
                if col_name in ["No", "Sentimen Lexicon", "Sentimen LLM", "Ground Truth"]:
                    cell.alignment = align_center
                else:
                    cell.alignment = align_left
                    
                # Apply conditional formatting for sentiment columns
                if col_name in ["Sentimen Lexicon", "Sentimen LLM", "Ground Truth"]:
                    val_lower = str(cell.value or "").strip().lower()
                    if col_name == "Ground Truth":
                        # Check mismatch between Sentimen Lexicon and Sentimen LLM in this row
                        # Row index in df_export is row_idx - 8
                        df_row_idx = row_idx - 8
                        lex_val = str(df_export.iloc[df_row_idx].get("Sentimen Lexicon", "")).strip().lower()
                        llm_val = str(df_export.iloc[df_row_idx].get("Sentimen LLM", "")).strip().lower()
                        
                        if lex_val != llm_val:
                            cell.fill = mismatch_fill
                            cell.font = mismatch_font
                        else:
                            if val_lower == "positif":
                                cell.fill = positif_fill
                                cell.font = positif_font
                            elif val_lower == "negatif":
                                cell.fill = negatif_fill
                                cell.font = negatif_font
                            elif val_lower == "netral":
                                cell.fill = netral_fill
                                cell.font = netral_font
                    else:
                        if val_lower == "positif":
                            cell.fill = positif_fill
                            cell.font = positif_font
                        elif val_lower == "negatif":
                            cell.fill = negatif_fill
                            cell.font = negatif_font
                        elif val_lower == "netral":
                            cell.fill = netral_fill
                            cell.font = netral_font
                
                # Estimate necessary row height dynamically by analyzing wrapped lines
                val = str(cell.value or "")
                val_lines = val.split('\n')
                width = col_widths.get(col_name, 15)
                # Count wrapped lines for this cell
                lines_in_cell = sum(max(1, int(np.ceil(len(l) / width))) for l in val_lines)
                max_lines = max(max_lines, lines_in_cell)
            
            # Set dynamic height: 14pt per line + 12pt padding (min height 20pt)
            worksheet.row_dimensions[row_idx].height = max(20, max_lines * 14 + 12)
            
        # Set Column Widths
        for col_idx, col_name in enumerate(df_export.columns, start=1):
            col_letter = get_column_letter(col_idx)
            width = col_widths.get(col_name, 15)
            worksheet.column_dimensions[col_letter].width = width
            
        # Set Header Row Heights
        worksheet.row_dimensions[1].height = 24
        worksheet.row_dimensions[2].height = 18
        worksheet.row_dimensions[3].height = 18
        worksheet.row_dimensions[4].height = 18 # Model LLM
        worksheet.row_dimensions[5].height = 18 # Language
        worksheet.row_dimensions[6].height = 12 # Empty spacer row
        worksheet.row_dimensions[7].height = 28 # Table header
            
    return output.getvalue()

# Helper to convert DataFrame to a PowerPoint presentation (.pptx) report in memory
def convert_df_to_pptx(df, video_title, video_url, analysis_mode, llm_model):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    
    prs = Presentation()
    # Set 16:9 widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # --- Color Palette ---
    CLR_PRIMARY  = RGBColor(30, 58, 138)   # Deep navy blue
    CLR_ACCENT   = RGBColor(59, 130, 246)  # Bright blue accent
    CLR_DARK     = RGBColor(30, 41, 59)    # Slate-900 for body text
    CLR_SUBTITLE = RGBColor(71, 85, 105)   # Slate-600
    CLR_LIGHT_BG = RGBColor(241, 245, 249) # Slate-100 (placeholder bg)
    CLR_BORDER   = RGBColor(203, 213, 225) # Slate-300 (placeholder border)
    CLR_PLACEHOLDER_TEXT = RGBColor(148, 163, 184) # Slate-400
    CLR_WHITE    = RGBColor(255, 255, 255)
    CLR_GREEN    = RGBColor(22, 163, 74)   # Green-600
    CLR_ORANGE   = RGBColor(234, 88, 12)   # Orange-600
    
    FONT_TITLE  = "Segoe UI"
    FONT_BODY   = "Segoe UI"
    
    # =============================================
    # Helper: Add a styled text bullet to text frame
    # =============================================
    def add_bullet(tf, text, level=0, bold=False, font_size=14, color=None, space_after=8):
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(space_after)
        p.space_before = Pt(2)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color or CLR_DARK
        run.font.bold = bold
        return p
    
    # =============================================
    # Helper: Create a content slide with accent bar, title, and image placeholder
    # =============================================
    def make_content_slide(title_text, img_placeholder_label=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # White background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = CLR_WHITE
        
        # Top accent bar (full width, 6px tall)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CLR_ACCENT
        bar.line.fill.background()
        
        # Left accent strip (vertical, 4px wide)
        vbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.35), Pt(5), Inches(0.55))
        vbar.fill.solid()
        vbar.fill.fore_color.rgb = CLR_ACCENT
        vbar.line.fill.background()
        
        # Title text box
        txBox = slide.shapes.add_textbox(Inches(0.9), Inches(0.3), Inches(7.5), Inches(0.65))
        tf_title = txBox.text_frame
        tf_title.word_wrap = True
        p = tf_title.paragraphs[0]
        p.text = title_text
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = FONT_TITLE
        run.font.size = Pt(26)
        run.font.color.rgb = CLR_PRIMARY
        run.font.bold = True
        
        # Separator line below title
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(4.8), Pt(2))
        sep.fill.solid()
        sep.fill.fore_color.rgb = CLR_ACCENT
        sep.line.fill.background()
        
        # Content text box (left column)
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(5.0), Inches(5.6))
        tf = txBox2.text_frame
        tf.word_wrap = True
        # Clear default paragraph
        tf.paragraphs[0].text = ""
        
        # Image placeholder (right column) — rounded rectangle
        if img_placeholder_label:
            ph_left = Inches(6.1)
            ph_top = Inches(1.35)
            ph_width = Inches(6.6)
            ph_height = Inches(5.6)
            
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ph_left, ph_top, ph_width, ph_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = CLR_LIGHT_BG
            shape.line.color.rgb = CLR_BORDER
            shape.line.width = Pt(1.5)
            shape.line.dash_style = 2  # Dashed
            
            # Placeholder label text
            ph_tf = shape.text_frame
            ph_tf.word_wrap = True
            ph_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            ph_p = ph_tf.paragraphs[0]
            ph_p.space_before = Pt(100)
            ph_run = ph_p.add_run()
            ph_run.text = f"📷\n{img_placeholder_label}"
            ph_run.font.name = FONT_BODY
            ph_run.font.size = Pt(12)
            ph_run.font.color.rgb = CLR_PLACEHOLDER_TEXT
            ph_run.font.bold = True
        
        # Slide number at bottom-right
        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.35))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.alignment = PP_ALIGN.RIGHT
        num_run = num_p.add_run()
        num_run.text = str(len(prs.slides))
        num_run.font.name = FONT_BODY
        num_run.font.size = Pt(10)
        num_run.font.color.rgb = CLR_SUBTITLE
        
        return tf
    
    # ====================================================
    # SLIDE 1 — Title Slide (Custom full-bleed design)
    # ====================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg1 = slide1.background
    fill1 = bg1.fill
    fill1.solid()
    fill1.fore_color.rgb = CLR_PRIMARY
    
    # Big brand accent block at left
    accent_block = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.6), Inches(7.5))
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = CLR_ACCENT
    accent_block.line.fill.background()
    
    # Course Label
    t_course = slide1.shapes.add_textbox(Inches(1.2), Inches(0.8), Inches(11.0), Inches(0.4))
    tf_course = t_course.text_frame
    p_course = tf_course.paragraphs[0]
    r_course = p_course.add_run()
    r_course.text = "TUGAS BESAR MATAKULIAH: STKI"
    r_course.font.name = FONT_TITLE
    r_course.font.size = Pt(14)
    r_course.font.color.rgb = RGBColor(191, 219, 254)  # Blue-200
    r_course.font.bold = True
    
    # Title
    t1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(11.0), Inches(2.2))
    tf1 = t1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "ANALISIS SENTIMEN KOMENTAR YOUTUBE\nDENGAN MODEL LLM & LEXICON"
    r1.font.name = FONT_TITLE
    r1.font.size = Pt(36)
    r1.font.color.rgb = CLR_WHITE
    r1.font.bold = True
    
    # Subtitle / Mode Comparison
    t_sub = slide1.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(11.0), Inches(0.6))
    tf_sub = t_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    r_sub = p_sub.add_run()
    r_sub.text = "Komparasi Performa: Konteks Global vs Konteks ke Video"
    r_sub.font.name = FONT_BODY
    r_sub.font.size = Pt(20)
    r_sub.font.color.rgb = CLR_ACCENT
    r_sub.font.bold = True
    
    # Student Info
    t_profile = slide1.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(11.0), Inches(0.8))
    tf_profile = t_profile.text_frame
    p_profile = tf_profile.paragraphs[0]
    r_profile = p_profile.add_run()
    r_profile.text = "Oleh:\nLuthfi Shidqi Habibulloh  (NPM: 96.2023.1.07702)"
    r_profile.font.name = FONT_BODY
    r_profile.font.size = Pt(16)
    r_profile.font.color.rgb = CLR_WHITE
    
    # Technical Metadata
    meta_lines = [
        f"Target Video :  {video_title}",
        f"Model LLM    :  {llm_model} (NVIDIA NIM)",
        f"Metode       :  Sastrawi + InSet Lexicon & VADER"
    ]
    t_meta = slide1.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(11.0), Inches(1.5))
    tf_meta = t_meta.text_frame
    tf_meta.word_wrap = True
    for i, line in enumerate(meta_lines):
        if i == 0:
            p = tf_meta.paragraphs[0]
        else:
            p = tf_meta.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = line
        r.font.name = FONT_BODY
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor(148, 163, 184)  # Slate-400

    # --- Pre-calculate all sentiment metrics for slides ---
    lex_pos = len(df[df["Lexicon Sentiment"].str.lower().str.strip() == "positif"])
    lex_neg = len(df[df["Lexicon Sentiment"].str.lower().str.strip() == "negatif"])
    lex_net = len(df[df["Lexicon Sentiment"].str.lower().str.strip() == "netral"])
    llm_pos = len(df[df["LLM Sentiment"].str.lower().str.strip() == "positif"])
    llm_neg = len(df[df["LLM Sentiment"].str.lower().str.strip() == "negatif"])
    llm_net = len(df[df["LLM Sentiment"].str.lower().str.strip() == "netral"])
    
    df_eval = df.dropna(subset=["Ground Truth"]).copy()
    df_eval = df_eval[df_eval["Ground Truth"].astype(str).str.strip().str.lower().isin(["positif", "negatif", "netral"])]
    
    lex_acc = 0.0
    llm_acc = 0.0
    has_eval = len(df_eval) > 0
    if has_eval:
        y_true = df_eval["Ground Truth"].str.strip().str.lower()
        y_lexicon = df_eval["Lexicon Sentiment"].str.strip().str.lower()
        y_llm = df_eval["LLM Sentiment"].str.strip().str.lower()
        lex_acc = accuracy_score(y_true, y_lexicon) if len(y_true) > 0 else 0
        llm_acc = accuracy_score(y_true, y_llm) if len(y_true) > 0 else 0
    
    # ====================================================
    # SLIDE 2 — Ringkasan Dataset Video Aktif (Dataset di Awal)
    # ====================================================
    tf = make_content_slide("Ringkasan Dataset Video", "Screenshot\nDashboard Utama")
    add_bullet(tf, "Ini dia video target yang mau saya demonstrasikan hari ini:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, f"Judul Video:  {video_title}", font_size=13, space_after=8)
    add_bullet(tf, f"URL Video:  {video_url}", font_size=13, space_after=8)
    add_bullet(tf, f"Total Data:  {len(df)} komentar asli netizen yang berhasil di-scrape", font_size=13, space_after=14)
    
    lang_info = "Indonesia"
    if "Language" in df.columns:
        lang_counts = df["Language"].value_counts()
        lang_parts = []
        for lang, count in lang_counts.items():
            pct = (count / len(df) * 100) if len(df) > 0 else 0
            name = "Indonesia" if str(lang).strip().lower() == "id" else ("Inggris" if str(lang).strip().lower() == "en" else str(lang).upper())
            lang_parts.append(f"{name} ({pct:.1f}%)")
        lang_info = ",  ".join(lang_parts)
    add_bullet(tf, "Detail Pengaturan:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, f"Bahasa Komentar:  {lang_info}", font_size=13, space_after=8)
    add_bullet(tf, f"Mode Analisis Aktif:  {analysis_mode}", font_size=13)
    
    # ====================================================
    # SLIDE 3 — Pendekatan Hibrida: Lexicon vs LLM
    # ====================================================
    tf = make_content_slide("Pendekatan Analisis: Lexicon vs LLM", "Tabel Perbandingan\nHasil Komentar")
    add_bullet(tf, "Untuk video tadi, saya langsung mengadu 2 metode analitik sekaligus:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Lexicon-Based  →  Deteksi instan pakai kamus kata (InSet & VADER)", font_size=13, space_after=8)
    add_bullet(tf, "LLM-Based  →  Model AI NVIDIA NIM buat pahami konteks kalimat", font_size=13, space_after=14)
    add_bullet(tf, "Keuntungan digabung?", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Kamus itu super kencang, tapi LLM pinter baca sarkasme & slang gaul", font_size=13, space_after=8)
    add_bullet(tf, "Perbandingan keduanya membuat saya mendapatkan kesimpulan yang objektif", font_size=13)
    
    # ====================================================
    # SLIDE 4 — Lexicon Bahasa Indonesia (Stemming + InSet)
    # ====================================================
    tf = make_content_slide("Lexicon Bahasa Indonesia", "Diagram / Screenshot\nHasil Lexicon Indonesia")
    add_bullet(tf, "Cara sistem olah komentar bahasa Indonesia di video ini:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Pemisahan komentar bahasa Indonesia dilakukan otomatis", font_size=13, space_after=8)
    add_bullet(tf, "Slang dibersihkan dulu (contoh: 'yg' -> 'yang', 'gk' -> 'tidak')", font_size=13, space_after=8)
    add_bullet(tf, "Kata berimbuhan diubah ke kata dasar pakai Stemmer Sastrawi", font_size=13, space_after=14)
    add_bullet(tf, "Pencocokan Kamus InSet:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Tiap kata dicocokkan ke kamus buat dicari skor positif/negatifnya", font_size=13, space_after=8)
    add_bullet(tf, "Skor akhir dijumlahkan buat tentuin sentimen komentarnya", font_size=13)
 
    # ====================================================
    # SLIDE 5 — Fitur Kustom: Pengelolaan Kamus Slang
    # ====================================================
    tf = make_content_slide("Pengelolaan Kamus Slang Kustom", "Tabel Kamus Slang\ndi Dashboard")
    add_bullet(tf, "Agar tebakan Kamus makin akurat, saya membuat editor Kamus Slang kustom:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Bisa tambah, edit, atau hapus padanan kata gaul sesuka hati di sidebar", font_size=13, space_after=8)
    add_bullet(tf, "Membantu menaikkan akurasi Lexicon secara drastis dari waktu ke waktu", font_size=13, space_after=14)
    add_bullet(tf, "Kelebihan Fitur Ini:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Sangat fleksibel untuk bahasa kekinian atau singkatan khas komunitas", font_size=13, space_after=8)
    add_bullet(tf, "Perubahan kata slang langsung diterapkan secara instan di analisis berikutnya", font_size=13)
    
    # ====================================================
    # SLIDE 6 — Lexicon Bahasa Inggris (VADER)
    # ====================================================
    tf = make_content_slide("Lexicon Bahasa Inggris — VADER", "Grafik Sebaran\nSentimen Inggris")
    add_bullet(tf, "Bagaimana dengan komentar bahasa Inggris di video tadi?", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Saya menggunakan VADER Lexicon yang sesuai untuk bahasa media sosial", font_size=13, space_after=8)
    add_bullet(tf, "Hebatnya, nggak perlu stemming buat nemu kata dasar", font_size=13, space_after=8)
    add_bullet(tf, "Pinter baca emoji (😊), tulisan caps lock (ANGRY), dan tanda seru (!!!)", font_size=13, space_after=14)
    add_bullet(tf, "Aturan Klasifikasi:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Menghasilkan compound score dari rentang -1 sampai +1", font_size=13, space_after=8)
    add_bullet(tf, "Batas sentimen: ≥ 0.05 Positif, ≤ -0.05 Negatif, sisanya Netral", font_size=13)
    
    # ====================================================
    # SLIDE 7 — Model LLM & Ground Truth
    # ====================================================
    tf = make_content_slide("Model LLM & Ground Truth", "Grafik Performa\nAkurasi Model")
    add_bullet(tf, "Dua jenis model AI di belakang layar:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Model Utama (Llama 3.1 8B)  →  Buat analisis cepat ratusan komentar", font_size=13, space_after=8)
    add_bullet(tf, "DeepSeek V4 Pro  →  Dipakai buat otomatisasi isi Ground Truth", font_size=13, space_after=14)
    add_bullet(tf, "Biar demo lancar dan aman dari error:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, f"Model aktif saat ini:  {llm_model.split('/')[-1]}", font_size=13, space_after=8)
    add_bullet(tf, "Ada auto-fallback ke model cadangan kalau API limit (Error 429)", font_size=13)
    
    # ====================================================
    # SLIDE 8 — Mode Analisis: Konteks Global
    # ====================================================
    tf = make_content_slide("Mode Analisis: Konteks Global", "Grafik Distribusi\nSentimen Global")
    add_bullet(tf, "Pertama, saya menguji video ini menggunakan Mode Konteks Global:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Evaluasi sentimen dilakukan mandiri per komentar (standalone)", font_size=13, space_after=8)
    add_bullet(tf, "Penilaian murni dari kalimat yang tertulis saja", font_size=13, space_after=8)
    add_bullet(tf, "Tidak mempertimbangkan isi atau topik video sama sekali", font_size=13, space_after=14)
    add_bullet(tf, "Kapan saya menggunakan mode ini?", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Pas mau tahu reaksi umum penonton tanpa butuh info isi video", font_size=13, space_after=8)
    add_bullet(tf, "Bagus buat membandingkan sentimen antar video secara global", font_size=13)
    
    # ====================================================
    # SLIDE 9 — Mode Konteks ke Video
    # ====================================================
    tf = make_content_slide("Mode Analisis: Konteks ke Video", "Screenshot Metadata\n& Transkrip Video")
    add_bullet(tf, "Kedua, saya menguji video yang sama menggunakan Mode Konteks ke Video:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Sistem kirim judul, deskripsi, tags, & transkrip video sebagai referensi LLM", font_size=13, space_after=8)
    add_bullet(tf, "Proses penarikan data pakai youtube-transcript-api secara instan", font_size=13, space_after=14)
    add_bullet(tf, "Efek ke hasil analisis pada dataset ini:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Model LLM bandingkan makna komentar langsung dengan isi video", font_size=13, space_after=8)
    add_bullet(tf, "Komentar spam atau out-of-context otomatis digeser ke Netral", font_size=13, space_after=8)
    add_bullet(tf, "Sentimen yang didapat jauh lebih presisi & sesuai konteks video", font_size=13)
 
    # ====================================================
    # SLIDE 10 — Hasil Sentimen & Akurasi
    # ====================================================
    tf = make_content_slide("Hasil Sentimen & Akurasi", "Pie / Bar Chart\nHasil Sentimen")
    add_bullet(tf, "Ini dia sebaran hasil tebakan model di video ini:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, f"Hasil Lexicon:  {lex_pos} Positif   |   {lex_neg} Negatif   |   {lex_net} Netral", font_size=13, space_after=8)
    add_bullet(tf, f"Hasil LLM:  {llm_pos} Positif   |   {llm_neg} Negatif   |   {llm_net} Netral", font_size=13, space_after=14)
    
    if has_eval:
        add_bullet(tf, "Skor Akurasi (vs Ground Truth):", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
        add_bullet(tf, f"Akurasi Lexicon:  {lex_acc*100:.1f}%", font_size=13, space_after=8)
        add_bullet(tf, f"Akurasi LLM:  {llm_acc*100:.1f}%", font_size=13, color=CLR_GREEN)
    else:
        add_bullet(tf, "⚠ Hasil akurasi belum muncul karena Ground Truth belum diisi", font_size=13, color=CLR_ORANGE)
 
    # ====================================================
    # SLIDE 11 — Pemodelan Topik (Topic Modeling)
    # ====================================================
    tf = make_content_slide("Pemodelan Topik (Topic Modeling)", "Pembagian Topik\ndi Tab Visualisasi")
    add_bullet(tf, "Selain sentimen, sistem memetakan tema obrolan di video ini:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Menggunakan algoritma pembobotan kata TF-IDF", font_size=13, space_after=8)
    add_bullet(tf, "Pengelompokan otomatis komentar pakai klusterisasi K-Means", font_size=13, space_after=14)
    add_bullet(tf, "Manfaat Analitik Ini:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Secara otomatis membagi komentar ke dalam 3 topik utama", font_size=13, space_after=8)
    add_bullet(tf, "Membantu saya mengetahui apa saja sub-topik yang paling banyak didebatkan", font_size=13)
 
    # ====================================================
    # SLIDE 12 — Analisis Perbandingan Global
    # ====================================================
    tf = make_content_slide("Analisis Perbandingan Global", "Grafik Perbandingan\nAntar Video")
    add_bullet(tf, "Terakhir, saya dapat membandingkan video ini dengan riwayat video lainnya:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Menggabungkan riwayat riil lokal maupun data dari cloud (GSheets)", font_size=13, space_after=8)
    add_bullet(tf, "Menampilkan grafik akurasi jangka panjang untuk melihat model terbaik", font_size=13, space_after=14)
    add_bullet(tf, "Fleksibilitas Filter:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Tersedia filter perbandingan khusus berdasarkan mode analisis", font_size=13, space_after=8)
    add_bullet(tf, "Bisa memilih / men-toggle video mana saja yang ingin disertakan", font_size=13)
    
    # ====================================================
    # SLIDE 13 — Kesimpulan & Temuan Utama
    # ====================================================
    tf = make_content_slide("Kesimpulan & Temuan Utama", "Bagan Kesimpulan\n/ Rekomendasi")
    add_bullet(tf, "Temuan penting dari uji coba saya:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=14)
    add_bullet(tf, "Metode Hibrida sukses bikin hasil tebakan saling melengkapi", font_size=13, space_after=8)
    add_bullet(tf, "Mode Konteks Video ampuh memangkas komentar nyasar & spam", font_size=13, space_after=8)
    add_bullet(tf, "DeepSeek V4 Pro terbukti menghasilkan Ground Truth yang masuk akal", font_size=13, space_after=14)
    add_bullet(tf, "Saran untuk penggunaan selanjutnya:", font_size=15, bold=True, color=CLR_PRIMARY, space_after=10)
    add_bullet(tf, "Pakai selalu Mode Konteks Video biar dapet hasil yang paling relevan", font_size=13, space_after=8)
    add_bullet(tf, "Jangan lupa review Ground Truth biar evaluasi metrik makin valid", font_size=13)
    # Save presentation to memory stream
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output.getvalue()

# Helper to convert DataFrame to a beautifully styled landscape A4 PDF report in memory
def convert_df_to_pdf(df, video_title, video_url):
    output = io.BytesIO()
    
    # Setup landscape A4 document
    doc = SimpleDocTemplate(
        output,
        pagesize=landscape(A4),
        leftMargin=30,
        rightMargin=30,
        topMargin=30,
        bottomMargin=30
    )
    
    story = []
    styles = getSampleStyleSheet()
    
    # Custom heading & metadata styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=16,
        textColor=colors.HexColor('#1F4E79'),
        spaceAfter=12
    )
    
    meta_style = ParagraphStyle(
        'MetaText',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9,
        textColor=colors.HexColor('#1E293B'),
        spaceAfter=4
    )
    
    meta_link_style = ParagraphStyle(
        'MetaLink',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        textColor=colors.HexColor('#2563EB'),
        spaceAfter=8
    )
    
    # Append report title & metadata headers
    story.append(Paragraph("SEMANTIKA - Laporan Analisis Sentimen Komentar YouTube", title_style))
    story.append(Paragraph(f"Judul Video: {video_title}", meta_style))
    story.append(Paragraph(f"Link Video: <font color='#2563EB'><u>{video_url}</u></font>", meta_link_style))
    story.append(Paragraph(f"Model LLM: {st.session_state.llm_model}", meta_style))
    story.append(Paragraph(f"Mode Analisis: {st.session_state.analysis_mode}", meta_style))
    lang_label = "Inggris (EN)" if st.session_state.detected_lang == "en" else "Indonesia (ID)"
    story.append(Paragraph(f"Bahasa Terdeteksi: {lang_label}", meta_style))
    story.append(Spacer(1, 15))
    
    # Column specific table styles (with wrap_text and alignment)
    th_style = ParagraphStyle(
        'TableHeader',
        fontName='Helvetica-Bold',
        fontSize=9,
        textColor=colors.white,
        alignment=1 # Centered
    )
    
    td_center_style = ParagraphStyle(
        'TableCellCenter',
        fontName='Helvetica',
        fontSize=8,
        textColor=colors.HexColor('#0F172A'),
        alignment=1 # Centered
    )
    
    td_left_style = ParagraphStyle(
        'TableCellLeft',
        fontName='Helvetica',
        fontSize=8,
        textColor=colors.HexColor('#0F172A'),
        alignment=0 # Left-aligned
    )
    
    # Table headers row
    headers = [
        Paragraph("<b>No</b>", th_style),
        Paragraph("<b>Penulis</b>", th_style),
        Paragraph("<b>Komentar Asli</b>", th_style),
        Paragraph("<b>Komentar Bersih (Stemmed)</b>", th_style),
        Paragraph("<b>Sentimen Lexicon</b>", th_style),
        Paragraph("<b>Sentimen LLM</b>", th_style),
        Paragraph("<b>Alasan LLM</b>", th_style),
        Paragraph("<b>Ground Truth</b>", th_style)
    ]
    
    table_data = [headers]
    
    # Format and append rows as Paragraph cells to allow word wrapping
    # Helper to truncate text to prevent row height exceeding page size and causing LayoutError.
    # Also escapes HTML characters to prevent ReportLab XML parser crashes on comments like '<3' or '&'.
    def cell_text(val, max_len=300):
        import html
        if val is None or pd.isna(val):
            return ""
        val_str = html.escape(str(val).strip())
        if len(val_str) > max_len:
            return val_str[:max_len] + "..."
        return val_str

    for idx, row in df.iterrows():
        no_p = Paragraph(str(idx + 1), td_center_style)
        author_p = Paragraph(cell_text(row.get("Author", ""), max_len=40), td_left_style)
        orig_p = Paragraph(cell_text(row.get("Original Comment", ""), max_len=300), td_left_style)
        clean_p = Paragraph(cell_text(row.get("Cleaned Comment", ""), max_len=300), td_left_style)
        lex_p = Paragraph(cell_text(row.get("Lexicon Sentiment", "")).capitalize(), td_center_style)
        llm_p = Paragraph(cell_text(row.get("LLM Sentiment", "")).capitalize(), td_center_style)
        reason_p = Paragraph(cell_text(row.get("LLM Reason", ""), max_len=300), td_left_style)
        
        gt_val = row.get("Ground Truth", "")
        gt_p = Paragraph(cell_text(gt_val).capitalize() if gt_val else "-", td_center_style)
        
        table_data.append([no_p, author_p, orig_p, clean_p, lex_p, llm_p, reason_p, gt_p])
        
    # Printable landscape A4 width is 781.89 points. Sum of columns: 780 pt.
    col_widths = [25, 75, 180, 120, 60, 60, 190, 70]
    
    # Create Table object
    t = Table(table_data, colWidths=col_widths, repeatRows=1)
    
    # Add table styling (background fills, margins, lines)
    t_style = [
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1F4E79')), # Header fill Slate Blue
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('TOPPADDING', (0,0), (-1,0), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')), # Light gray gridlines
    ]
    
    # Apply zebra-striping rows and conditional background colors for sentiment columns
    for r in range(1, len(table_data)):
        bg_color = colors.HexColor('#F8FAFC') if r % 2 == 0 else colors.white
        t_style.append(('BACKGROUND', (0, r), (-1, r), bg_color))
        t_style.append(('TOPPADDING', (0, r), (-1, r), 5))
        t_style.append(('BOTTOMPADDING', (0, r), (-1, r), 5))
        
        # Get matching DataFrame row (header is row 0 in table_data)
        row = df.iloc[r - 1]
        
        # Lexicon Sentiment (Column index 4)
        lex_val = str(row.get("Lexicon Sentiment", "")).strip().lower()
        if lex_val == "positif":
            t_style.append(('BACKGROUND', (4, r), (4, r), colors.HexColor('#C6EFCE')))
        elif lex_val == "negatif":
            t_style.append(('BACKGROUND', (4, r), (4, r), colors.HexColor('#FFC7CE')))
        elif lex_val == "netral":
            t_style.append(('BACKGROUND', (4, r), (4, r), colors.HexColor('#E2E3E5')))
            
        # LLM Sentiment (Column index 5)
        llm_val = str(row.get("LLM Sentiment", "")).strip().lower()
        if llm_val == "positif":
            t_style.append(('BACKGROUND', (5, r), (5, r), colors.HexColor('#C6EFCE')))
        elif llm_val == "negatif":
            t_style.append(('BACKGROUND', (5, r), (5, r), colors.HexColor('#FFC7CE')))
        elif llm_val == "netral":
            t_style.append(('BACKGROUND', (5, r), (5, r), colors.HexColor('#E2E3E5')))
            
        # Ground Truth (Column index 7)
        gt_val = str(row.get("Ground Truth", "")).strip().lower()
        if lex_val != llm_val:
            t_style.append(('BACKGROUND', (7, r), (7, r), colors.HexColor('#FFE699')))
        else:
            if gt_val == "positif":
                t_style.append(('BACKGROUND', (7, r), (7, r), colors.HexColor('#C6EFCE')))
            elif gt_val == "negatif":
                t_style.append(('BACKGROUND', (7, r), (7, r), colors.HexColor('#FFC7CE')))
            elif gt_val == "netral":
                t_style.append(('BACKGROUND', (7, r), (7, r), colors.HexColor('#E2E3E5')))
        
    t.setStyle(TableStyle(t_style))
    story.append(t)
    
    # Compile document
    doc.build(story)
    return output.getvalue()

# Sidebar Navigation Menu
st.sidebar.title(":material/explore: Navigasi Menu")
menu_selection = st.sidebar.radio(
    "Pilih Halaman:",
    options=["Analisis Video Tunggal", "Analisis Perbandingan Global", "Kelola Kamus Slang"],
    index=0
)
st.sidebar.markdown("---")

def load_all_gsheets_data():
    """
    Mengambil seluruh data riwayat sentimen dari worksheet Google Sheets 'Database_Sentimen'.
    """
    cols = ["Video ID", "Video Title", "Video URL", "Comment ID", "Author", "Original Comment", "Cleaned Comment", "Lexicon Sentiment", "Lexicon Score", "LLM Sentiment", "LLM Reason", "LLM Model", "Language", "Analysis Mode", "Ground Truth"]
    if APP_MODE != "production" or conn is None:
        return pd.DataFrame(columns=cols)
        
    try:
        df = conn.read(worksheet="Database_Sentimen", ttl="0")
        if df is not None and not df.empty:
            for col in cols:
                if col not in df.columns:
                    df[col] = ""
            return df[cols]
    except Exception:
        pass
    return pd.DataFrame(columns=cols)

def save_to_gsheets(df_to_save):
    """
    Menyimpan DataFrame secara penuh kembali ke worksheet Google Sheets 'Database_Sentimen'.
    """
    if APP_MODE != "production" or conn is None:
        return
        
    try:
        df_clean = df_to_save.copy()
        for col in df_clean.columns:
            df_clean[col] = df_clean[col].fillna("").astype(str)
        conn.update(worksheet="Database_Sentimen", data=df_clean)
    except Exception as e:
        st.error(f"Gagal sinkronisasi ke Google Sheets: {e}")

def sync_video_to_gsheets(video_id, video_title, video_url, df_video):
    """
    Menambahkan/memperbarui data hasil analisis video tertentu ke master Google Sheets.
    """
    if APP_MODE != "production" or conn is None:
        return
        
    df_all = load_all_gsheets_data()
    
    # Hapus baris yang memiliki Video ID dan Mode Analisis yang sama agar tidak duplikat
    if not df_all.empty:
        analysis_mode = st.session_state.analysis_mode
        if "Analysis Mode" in df_video.columns and len(df_video) > 0:
            analysis_mode = df_video["Analysis Mode"].iloc[0]
        df_all = df_all[~((df_all["Video ID"] == video_id) & (df_all["Analysis Mode"] == analysis_mode))]
        
    new_rows = []
    for _, row in df_video.iterrows():
        new_rows.append({
            "Video ID": video_id,
            "Video Title": video_title,
            "Video URL": video_url,
            "Comment ID": row["Comment ID"] if "Comment ID" in row else row.get("comment_id", ""),
            "Author": row["Author"] if "Author" in row else row.get("author", ""),
            "Original Comment": row["Original Comment"] if "Original Comment" in row else row.get("original_comment", ""),
            "Cleaned Comment": row["Cleaned Comment"] if "Cleaned Comment" in row else row.get("cleaned_comment", ""),
            "Lexicon Sentiment": row["Lexicon Sentiment"] if "Lexicon Sentiment" in row else row.get("lexicon_sentiment", ""),
            "Lexicon Score": row["Lexicon Score"] if "Lexicon Score" in row else row.get("lexicon_score", ""),
            "LLM Sentiment": row["LLM Sentiment"] if "LLM Sentiment" in row else row.get("llm_sentiment", ""),
            "LLM Reason": row["LLM Reason"] if "LLM Reason" in row else row.get("llm_reason", ""),
            "LLM Model": row["LLM Model"] if "LLM Model" in row else st.session_state.llm_model,
            "Language": row["Language"] if "Language" in row else st.session_state.detected_lang,
            "Analysis Mode": row["Analysis Mode"] if "Analysis Mode" in row else st.session_state.analysis_mode,
            "Ground Truth": row.get("Ground Truth", "")
        })
    df_new = pd.DataFrame(new_rows)
    
    df_combined = pd.concat([df_all, df_new], ignore_index=True)
    save_to_gsheets(df_combined)

def get_existing_ground_truths():
    """
    Mengambil ground truth yang sudah diisi sebelumnya dari file hasil aktif (OUTPUT_FILE),
    seluruh berkas riwayat di folder history, dan database Google Sheets (jika dalam mode production).
    """
    gts = {}
    
    # 1. Baca dari sentiment_results.csv jika ada
    if os.path.exists(OUTPUT_FILE):
        try:
            df_old = pd.read_csv(OUTPUT_FILE)
            if "Comment ID" in df_old.columns and "Ground Truth" in df_old.columns:
                df_old = df_old.dropna(subset=["Comment ID"])
                for _, row in df_old.iterrows():
                    cid = str(row["Comment ID"]).strip()
                    gt = str(row["Ground Truth"]).strip() if pd.notna(row["Ground Truth"]) else ""
                    if cid and gt:
                        gts[cid] = gt
        except Exception:
            pass

    # 2. Baca dari folder history
    if os.path.exists(HISTORY_DIR):
        try:
            for f in os.listdir(HISTORY_DIR):
                if f.endswith(".csv"):
                    fpath = os.path.join(HISTORY_DIR, f)
                    df_hist = pd.read_csv(fpath)
                    if "Comment ID" in df_hist.columns and "Ground Truth" in df_hist.columns:
                        df_hist = df_hist.dropna(subset=["Comment ID"])
                        for _, row in df_hist.iterrows():
                            cid = str(row["Comment ID"]).strip()
                            gt = str(row["Ground Truth"]).strip() if pd.notna(row["Ground Truth"]) else ""
                            if cid and gt:
                                gts[cid] = gt
        except Exception:
            pass
            
    # 3. Baca dari Google Sheets (hanya di mode production)
    if APP_MODE == "production" and conn is not None:
        try:
            df_gs = load_all_gsheets_data()
            if not df_gs.empty and "Comment ID" in df_gs.columns and "Ground Truth" in df_gs.columns:
                df_gs = df_gs.dropna(subset=["Comment ID"])
                for _, row in df_gs.iterrows():
                    cid = str(row["Comment ID"]).strip()
                    gt = str(row["Ground Truth"]).strip() if pd.notna(row["Ground Truth"]) else ""
                    if cid and gt:
                        gts[cid] = gt
        except Exception:
            pass
            
    return gts


def make_safe_filename(name):
    return re.sub(r'[\\/*?:"<>|]', "", name).strip()

def detect_language_from_title(title):
    """
    Detects if the video title is primarily Indonesian ('id') or English ('en') using NVIDIA NIM.
    """
    detector = LLMSentimentAnalyzer(model="meta/llama-3.1-8b-instruct")
    system_prompt = (
        "You are a language detection assistant. Detect whether the following YouTube video title is primarily "
        "in Indonesian (or Indonesian slang/slang) or English.\n"
        "Respond with only 'id' for Indonesian/slang or 'en' for English. Do not write any other words or characters."
    )
    user_prompt = f"Video Title: {title}"
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt}
    ]
    try:
        response = detector._call_nvidia_api(messages)
        res_clean = response.strip().lower()
        if "en" in res_clean:
            return "en"
        return "id"
    except Exception:
        return "id"

if menu_selection == "Analisis Video Tunggal":
    # Sidebar Config
    st.sidebar.title(":material/settings: Setelan SEMANTIKA")
    st.sidebar.markdown("---")
    
    data_source = st.sidebar.radio(
        "Sumber Data",
        options=["Ambil Video Baru (Scraping)", "Buka Riwayat Analisis (Lokal)"],
        help="Pilih apakah ingin melakukan penarikan data baru dari YouTube atau membuka riwayat yang sudah ada di folder history."
    )
    
    # Initialize placeholders to avoid reference errors
    url_input = ""
    limit_input = MAX_COMMENTS
    model_input = st.session_state.llm_model
    force_refresh = False
    btn_analyze = False
    
    if data_source == "Ambil Video Baru (Scraping)":
        url_input = st.sidebar.text_input(
            "URL Video YouTube / Shorts",
            key="youtube_url_widget",
            help="Masukkan URL video YouTube atau Shorts yang ingin dianalisis."
        )
        
        limit_input = st.sidebar.slider(
            "Jumlah Komentar Maksimal",
            min_value=10,
            max_value=1000,
            value=100,
            step=10,
            help="Batasi jumlah komentar yang akan ditarik. Default adalah 100 komentar."
        )
        
        options_list = [
            "meta/llama-3.1-8b-instruct", 
            "meta/llama-3.1-70b-instruct", 
            "meta/llama-3.3-70b-instruct",
            "deepseek-ai/deepseek-v4-flash",
            "deepseek-ai/deepseek-v4-pro"
        ]
        default_index = 0
        if st.session_state.llm_model in options_list:
            default_index = options_list.index(st.session_state.llm_model)
        
        def format_model_name(model_id):
            labels = {
                "deepseek-ai/deepseek-v4-pro": "deepseek-ai/deepseek-v4-pro (Tercanggih - DeepSeek V4 Pro)",
                "meta/llama-3.3-70b-instruct": "meta/llama-3.3-70b-instruct (Sangat Canggih - Llama 3.3)",
                "meta/llama-3.1-70b-instruct": "meta/llama-3.1-70b-instruct (Canggih - Llama 3.1 70B)",
                "deepseek-ai/deepseek-v4-flash": "deepseek-ai/deepseek-v4-flash (Cepat & Canggih)",
                "meta/llama-3.1-8b-instruct": "meta/llama-3.1-8b-instruct (Cepat / Ringan)"
            }
            return labels.get(model_id, model_id)

        model_input = st.sidebar.selectbox(
            "Model LLM NVIDIA",
            options=options_list,
            index=default_index,
            format_func=format_model_name,
            help="Pilih model NVIDIA NIM yang ingin digunakan untuk klasifikasi."
        )
        st.session_state.analysis_mode = "Dual Mode"

        force_refresh = st.sidebar.checkbox(
            "Paksa Ambil Baru (Force Refresh)",
            value=False,
            help="Centang ini untuk mengabaikan riwayat lokal dan mengambil data baru dari YouTube & NVIDIA API."
        )
        
        btn_analyze = st.sidebar.button(":material/play_circle: Mulai Analisis Data", use_container_width=True)
    else:
        # Scan history files
        history_files = sorted([f for f in os.listdir(HISTORY_DIR) if f.endswith(".csv")])
        if not history_files:
            st.sidebar.warning("Tidak ditemukan riwayat analisis lokal di folder history.")
        else:
            selected_file = st.sidebar.selectbox(
                "Pilih Riwayat Analisis",
                options=history_files,
                help="Pilih file riwayat yang ingin ditampilkan."
            )
            btn_load_history = st.sidebar.button(":material/folder_open: Buka Riwayat", use_container_width=True)
            if btn_load_history and selected_file:
                history_path = os.path.join(HISTORY_DIR, selected_file)
                try:
                    df_loaded = pd.read_csv(history_path)
                    df_loaded = upgrade_dataframe_schema(df_loaded)
                    
                    # Extract video ID and Title from filename if format is "[id] Title.csv"
                    match = re.match(r"^\[([a-zA-Z0-9_-]+)\]\s*(.*)\.csv$", selected_file)
                    if match:
                        video_id = match.group(1)
                        video_title = re.sub(r"^\[(?:Konteks Global|Konteks ke Video)\]\s*", "", match.group(2)).strip()
                    else:
                        video_id = "unknown"
                        video_title = selected_file.replace(".csv", "")
                        
                    # Save as current output file
                    df_loaded.to_csv(OUTPUT_FILE, index=False, encoding="utf-8-sig")
                    
                    st.session_state.df = df_loaded
                    st.session_state.video_title = video_title
                    st.session_state.video_url = f"https://www.youtube.com/watch?v={video_id}" if video_id != "unknown" else ""
                    
                    if "LLM Model" in df_loaded.columns:
                        st.session_state.llm_model = str(df_loaded["LLM Model"].iloc[0])
                    if "Language" in df_loaded.columns:
                        st.session_state.detected_lang = str(df_loaded["Language"].iloc[0])
                    st.session_state.analysis_mode = "Dual Mode"
                        
                    st.session_state.video_context = None
                    if "Lexicon Time" in df_loaded.columns:
                        st.session_state.lexicon_time = float(df_loaded["Lexicon Time"].dropna().iloc[0]) if len(df_loaded["Lexicon Time"].dropna()) > 0 else len(df_loaded) * 0.0015
                    else:
                        st.session_state.lexicon_time = len(df_loaded) * 0.0015
                        
                    if "LLM Time Global" in df_loaded.columns:
                        st.session_state.llm_time_global = float(df_loaded["LLM Time Global"].dropna().iloc[0]) if len(df_loaded["LLM Time Global"].dropna()) > 0 else len(df_loaded) * 0.15
                    else:
                        st.session_state.llm_time_global = len(df_loaded) * 0.15
                        
                    if "LLM Time Video" in df_loaded.columns:
                        st.session_state.llm_time_video = float(df_loaded["LLM Time Video"].dropna().iloc[0]) if len(df_loaded["LLM Time Video"].dropna()) > 0 else len(df_loaded) * 0.15
                    else:
                        st.session_state.llm_time_video = len(df_loaded) * 0.15
                    st.sidebar.success("Riwayat berhasil dimuat!")
                    st.rerun()
                except Exception as e:
                    st.sidebar.error(f"Gagal memuat riwayat: {e}")
    
    if btn_analyze:
        st.session_state.loaded_history_file = ""
        if not url_input.strip():
            st.sidebar.error("Silakan masukkan URL YouTube terlebih dahulu!")
        else:
            video_id = extract_video_id(url_input)
            if not video_id:
                st.sidebar.error("Gagal mengekstrak Video ID dari URL!")
            else:
                video_title = get_video_title(url_input)
                safe_title = make_safe_filename(video_title)
                history_filename = f"[{video_id}] {safe_title}.csv"
                history_path = os.path.join(HISTORY_DIR, history_filename)
    
                if not force_refresh and os.path.exists(history_path):
                    st.sidebar.info("Hasil analisis ditemukan di riwayat lokal. Memuat...")
                    try:
                        df_loaded = pd.read_csv(history_path)
                        df_loaded = upgrade_dataframe_schema(df_loaded)
                        
                        if "LLM Model" in df_loaded.columns:
                            st.session_state.llm_model = str(df_loaded["LLM Model"].iloc[0])
                        else:
                            df_loaded["LLM Model"] = "meta/llama-3.1-8b-instruct"
                            st.session_state.llm_model = "meta/llama-3.1-8b-instruct"
                        if "Language" in df_loaded.columns:
                            st.session_state.detected_lang = str(df_loaded["Language"].iloc[0])
                        else:
                            df_loaded["Language"] = "id"
                            st.session_state.detected_lang = "id"
                            
                        st.session_state.analysis_mode = "Dual Mode"
                        df_loaded.to_csv(OUTPUT_FILE, index=False, encoding="utf-8-sig")
                        
                        st.session_state.df = df_loaded
                        st.session_state.video_title = video_title
                        st.session_state.video_url = url_input
                        st.session_state.video_context = None
                        if "Lexicon Time" in df_loaded.columns:
                            st.session_state.lexicon_time = float(df_loaded["Lexicon Time"].dropna().iloc[0]) if len(df_loaded["Lexicon Time"].dropna()) > 0 else len(df_loaded) * 0.0015
                        else:
                            st.session_state.lexicon_time = len(df_loaded) * 0.0015
                            
                        if "LLM Time Global" in df_loaded.columns:
                            st.session_state.llm_time_global = float(df_loaded["LLM Time Global"].dropna().iloc[0]) if len(df_loaded["LLM Time Global"].dropna()) > 0 else len(df_loaded) * 0.15
                        else:
                            st.session_state.llm_time_global = len(df_loaded) * 0.15
                            
                        if "LLM Time Video" in df_loaded.columns:
                            st.session_state.llm_time_video = float(df_loaded["LLM Time Video"].dropna().iloc[0]) if len(df_loaded["LLM Time Video"].dropna()) > 0 else len(df_loaded) * 0.15
                        else:
                            st.session_state.llm_time_video = len(df_loaded) * 0.15
                        st.rerun()
                    except Exception as e:
                        st.sidebar.error(f"Gagal memuat file riwayat: {e}")
                else:
                    with st.status("Menjalankan Analisis Sentimen...", expanded=True) as status:
                        # 1. Fetch Title & Context
                        status.write("Langkah 1/5: Mengambil informasi video YouTube...")
                        video_title = get_video_title(url_input)
                        safe_title = make_safe_filename(video_title)
                        history_filename = f"[{video_id}] {safe_title}.csv"
                        history_path = os.path.join(HISTORY_DIR, history_filename)
                        
                        status.write("   - Mengambil teks isi dan metadata video (konteks)...")
                        video_context = get_video_context(url_input)
                        
                        # Language Detection
                        status.write("   - Mendeteksi bahasa konten video...")
                        detected_lang = detect_language_from_title(video_title)
                        status.write(f"   - Bahasa terdeteksi: {detected_lang.upper()}")
                        
                        # 2. Fetch Comments
                        status.write("Langkah 2/5: Mengunduh komentar dari YouTube...")
                        comments = fetch_youtube_comments(url_input, limit=limit_input)
                        
                        if not comments:
                            status.update(label="Gagal mengambil komentar!", state="error", expanded=True)
                            st.error("Gagal mendapatkan komentar dari video ini.")
                        else:
                            status.write(f"   - Sukses mengunduh {len(comments)} komentar.")
                            
                            # 3. Analyze Lexicon
                            status.write("Langkah 3/5: Menjalankan pemrosesan & skoring Lexicon...")
                            import time
                            start_lexicon_time = time.time()
                            lexicon_analyzer = LexiconSentimentAnalyzer()
                            processed_comments = []
                            
                            for idx, c in enumerate(comments):
                                sentiment, score, cleaned_text, comment_lang = lexicon_analyzer.analyze_sentiment(c["text"], default_lang=detected_lang)
                                processed_comments.append({
                                    "comment_id": c["comment_id"],
                                    "author": c["author"],
                                    "original_comment": c["text"],
                                    "cleaned_comment": cleaned_text,
                                    "lexicon_sentiment": sentiment,
                                    "lexicon_score": score,
                                    "language": comment_lang,
                                    "likes": c.get("likes", 0),
                                    "time": c.get("time", ""),
                                    "time_parsed": c.get("time_parsed", 0.0)
                                })
                                if (idx + 1) % 10 == 0 or (idx + 1) == len(comments):
                                    status.write(f"   - Selesai memproses Lexicon: {idx + 1}/{len(comments)} komentar...")
                            st.session_state.lexicon_time = time.time() - start_lexicon_time
                            
                            # 4. Analyze LLM (Both Modes)
                            status.write(f"Langkah 4/5: Menghubungi NVIDIA NIM API ({model_input})...")
                            llm_analyzer = LLMSentimentAnalyzer(model=model_input)
                            batch_size = 20
                            num_batches = (len(comments) - 1) // batch_size + 1
                            
                            # Mode 1: Konteks Global
                            status.write("   - Memproses Mode Konteks Global...")
                            start_global_time = time.time()
                            llm_sentiment_global = {}
                            llm_reason_global = {}
                            try:
                                for batch_idx, i in enumerate(range(0, len(comments), batch_size)):
                                    batch = comments[i:i+batch_size]
                                    status.write(f"     [Global] Batch {batch_idx + 1}/{num_batches}...")
                                    batch_results = llm_analyzer.analyze_batch(batch, video_context=None)
                                    for r in batch_results:
                                        llm_sentiment_global[r["comment_id"]] = r["llm_sentiment"]
                                        llm_reason_global[r["comment_id"]] = r.get("llm_reason", "")
                            except Exception as e:
                                status.update(label="Gagal menghubungi API LLM (Konteks Global)!", state="error", expanded=True)
                                st.error(f"Terjadi kesalahan saat menghubungi API NVIDIA NIM: {e}.")
                                st.stop()
                            st.session_state.llm_time_global = time.time() - start_global_time
                            
                            # Mode 2: Konteks ke Video
                            status.write("   - Memproses Mode Konteks ke Video...")
                            start_video_time = time.time()
                            llm_sentiment_video = {}
                            llm_reason_video = {}
                            try:
                                for batch_idx, i in enumerate(range(0, len(comments), batch_size)):
                                    batch = comments[i:i+batch_size]
                                    status.write(f"     [Video] Batch {batch_idx + 1}/{num_batches}...")
                                    batch_results = llm_analyzer.analyze_batch(batch, video_context=video_context)
                                    for r in batch_results:
                                        llm_sentiment_video[r["comment_id"]] = r["llm_sentiment"]
                                        llm_reason_video[r["comment_id"]] = r.get("llm_reason", "")
                            except Exception as e:
                                status.update(label="Gagal menghubungi API LLM (Konteks ke Video)!", state="error", expanded=True)
                                st.error(f"Terjadi kesalahan saat menghubungi API NVIDIA NIM: {e}.")
                                st.stop()
                            st.session_state.llm_time_video = time.time() - start_video_time
                            
                            # Obtain DeepSeek V4 Pro sentiments for Ground Truth
                            ds_sentiment_map = {}
                            if model_input == "deepseek-ai/deepseek-v4-pro":
                                ds_sentiment_map = llm_sentiment_video
                            else:
                                ds_sentiment_map = {}
                                
                            # 5. Combine results and map existing ground truths
                            status.write("Langkah 5/5: Menyimpan berkas hasil...")
                            existing_gts = get_existing_ground_truths()
                            
                            final_data = []
                            for idx, c in enumerate(processed_comments):
                                cid = c["comment_id"]
                                gt = existing_gts.get(cid, "")
                                if not gt or str(gt).strip() == "":
                                    gt = ds_sentiment_map.get(cid, "")
                                final_data.append({
                                    "No": idx + 1,
                                    "Comment ID": cid,
                                    "Author": c["author"],
                                    "Original Comment": c["original_comment"],
                                    "Cleaned Comment": c["cleaned_comment"],
                                    "Likes": c["likes"],
                                    "Time Description": c["time"],
                                    "Timestamp": c["time_parsed"],
                                    "Lexicon Sentiment": c["lexicon_sentiment"],
                                    "Lexicon Score": c["lexicon_score"],
                                    "LLM Sentiment Global": llm_sentiment_global.get(cid, "netral"),
                                    "LLM Reason Global": llm_reason_global.get(cid, ""),
                                    "LLM Sentiment Video": llm_sentiment_video.get(cid, "netral"),
                                    "LLM Reason Video": llm_reason_video.get(cid, ""),
                                    "LLM Model": model_input,
                                    "Language": c["language"],
                                    "Ground Truth": gt,
                                    "Lexicon Time": st.session_state.lexicon_time,
                                    "LLM Time Global": st.session_state.llm_time_global,
                                    "LLM Time Video": st.session_state.llm_time_video
                                })
                                
                            df = pd.DataFrame(final_data)
                            df.to_csv(OUTPUT_FILE, index=False, encoding="utf-8-sig")
                            df.to_csv(history_path, index=False, encoding="utf-8-sig")
                            
                            if APP_MODE == "production":
                                sync_video_to_gsheets(video_id, video_title, url_input, df)
                            
                            st.session_state.df = df
                            st.session_state.video_title = video_title
                            st.session_state.video_url = url_input
                            st.session_state.video_context = video_context
                            st.session_state.llm_model = model_input
                            st.session_state.detected_lang = detected_lang
                            
                            status.update(label="Analisis sentimen berhasil diselesaikan!", state="complete", expanded=False)
                            st.rerun()
    
    # Sidebar: History Loading Section
    st.sidebar.markdown("---")
    st.sidebar.subheader(":material/history: Riwayat Analisis")
    
    history_files = []
    if os.path.exists(HISTORY_DIR):
        history_files = sorted(
            [f for f in os.listdir(HISTORY_DIR) if f.endswith(".csv")],
            key=lambda x: os.path.getmtime(os.path.join(HISTORY_DIR, x)),
            reverse=True
        )
    
    if history_files:
        def load_history_callback():
            selected = st.session_state.history_selector
            if selected != "-- Pilih untuk memuat --" and selected != st.session_state.loaded_history_file:
                history_path = os.path.join(HISTORY_DIR, selected)
                try:
                    df_loaded = pd.read_csv(history_path)
                    df_loaded = upgrade_dataframe_schema(df_loaded)
                    
                    if "LLM Model" in df_loaded.columns:
                        st.session_state.llm_model = str(df_loaded["LLM Model"].iloc[0])
                    else:
                        df_loaded["LLM Model"] = "meta/llama-3.1-8b-instruct"
                        st.session_state.llm_model = "meta/llama-3.1-8b-instruct"
                    if "Language" in df_loaded.columns:
                        st.session_state.detected_lang = str(df_loaded["Language"].iloc[0])
                    else:
                        df_loaded["Language"] = "id"
                        st.session_state.detected_lang = "id"
                    if "Analysis Mode" in df_loaded.columns:
                        st.session_state.analysis_mode = str(df_loaded["Analysis Mode"].iloc[0])
                    else:
                        df_loaded["Analysis Mode"] = "Konteks Global"
                        st.session_state.analysis_mode = "Konteks Global"
                    df_loaded.to_csv(OUTPUT_FILE, index=False, encoding="utf-8-sig")
                    
                    filename_clean = selected[:-4]
                    match = re.match(r"^\[(.*?)\] (.*)$", filename_clean)
                    if match:
                        vid_id = match.group(1)
                        vid_title = re.sub(r"^\[(?:Konteks Global|Konteks ke Video)\]\s*", "", match.group(2)).strip()
                        new_url = f"https://www.youtube.com/watch?v={vid_id}"
                    else:
                        new_url = YOUTUBE_VIDEO_URL
                        vid_title = filename_clean
                        
                    st.session_state.video_url = new_url
                    st.session_state.video_title = vid_title
                    st.session_state.df = df_loaded
                    st.session_state.youtube_url_widget = new_url  # Update sidebar widget safely in callback
                    st.session_state.loaded_history_file = selected
                    st.session_state.lexicon_time = None
                    st.session_state.llm_time = None
                    st.session_state.load_success_msg = "Berhasil memuat data riwayat!"
                except Exception as e:
                    st.session_state.load_error_msg = f"Gagal memuat: {e}"

        history_options = ["-- Pilih untuk memuat --"] + history_files
        selected_history = st.sidebar.selectbox(
            "Muat Hasil Sebelumnya",
            options=history_options,
            index=0,
            key="history_selector",
            on_change=load_history_callback,
            help="Muat hasil analisis secara instan dari lokal disk."
        )
        
        # Display messages from callback if any
        if "load_success_msg" in st.session_state and st.session_state.load_success_msg:
            st.sidebar.success(st.session_state.load_success_msg)
            st.session_state.load_success_msg = ""
        if "load_error_msg" in st.session_state and st.session_state.load_error_msg:
            st.sidebar.error(st.session_state.load_error_msg)
            st.session_state.load_error_msg = ""
    
        # Kelola Riwayat Expander
        with st.sidebar.expander(":material/delete: Kelola Riwayat"):
            st.markdown("<small>Centang riwayat yang ingin dihapus:</small>", unsafe_allow_html=True)
            select_all = st.checkbox("Pilih Semua", key="select_all_del")
            
            to_delete = []
            for h_file in history_files:
                clean_name = h_file.replace(".csv", "")
                match = re.match(r"^\[(.*?)\] (.*)$", clean_name)
                display_name = match.group(2) if match else clean_name
                if len(display_name) > 25:
                    display_name = display_name[:22] + "..."
                
                checked = st.checkbox(display_name, value=select_all, key=f"del_{h_file}")
                if checked:
                    to_delete.append(h_file)
                    
            col_del1, col_del2 = st.columns(2)
            with col_del1:
                if st.button("Hapus Terpilih", type="primary", use_container_width=True):
                    if to_delete:
                        for h_file in to_delete:
                            file_path = os.path.join(HISTORY_DIR, h_file)
                            if os.path.exists(file_path):
                                os.remove(file_path)
                        st.sidebar.success(f"Berhasil menghapus {len(to_delete)} riwayat!")
                        if st.session_state.loaded_history_file in to_delete:
                            st.session_state.df = None
                            st.session_state.loaded_history_file = ""
                            if os.path.exists(OUTPUT_FILE):
                                os.remove(OUTPUT_FILE)
                        st.rerun()
                    else:
                        st.sidebar.warning("Pilih riwayat dulu!")
            with col_del2:
                if st.button("Hapus Semua", use_container_width=True):
                    for h_file in history_files:
                        file_path = os.path.join(HISTORY_DIR, h_file)
                        if os.path.exists(file_path):
                            os.remove(file_path)
                    if os.path.exists(OUTPUT_FILE):
                        os.remove(OUTPUT_FILE)
                    st.session_state.df = None
                    st.session_state.loaded_history_file = ""
                    st.sidebar.success("Semua riwayat berhasil dihapus!")
                    st.rerun()
    else:
        st.sidebar.info("Belum ada riwayat analisis.")
else:
    # menu_selection == "Analisis Perbandingan Global"
    st.sidebar.title(":material/analytics: Perbandingan Global")
    st.sidebar.markdown("---")
    st.sidebar.info("Filter video dan toggle riwayat dikelola langsung pada panel di halaman utama.")

# Info Metodologi di Sidebar (Pojok Halaman)
st.sidebar.markdown("---")
methodology_md = """
### :material/info: Metodologi: Lexicon vs LLM

#### 1. Lexicon-based (Kamus Kata)
*   **Cara kerja:** Menjumlahkan skor/bobot sentimen kata demi kata berdasarkan kamus kosakata (*InSet* untuk ID / *VADER* untuk EN).
*   **Kelebihan:** Sangat cepat, transparan, dan tidak bergantung pada API eksternal.
*   **Kekurangan:** Tidak memahami konteks kalimat, sindiran (sarkasme), kata negasi (contoh: *"tidak jelek"* dideteksi negatif karena kata *"jelek"*), dan rentan salah jika ada kesalahan ejaan (typo) atau slang yang tidak terdaftar di kamus.

#### 2. LLM-based (Konteks AI / Semantik)
*   **Cara kerja:** Memahami keseluruhan kalimat secara utuh menggunakan kecerdasan buatan (NVIDIA NIM).
*   **Kelebihan:** Sangat pintar memahami konteks, sindiran, negasi, slang internet terbaru, singkatan ekstrim, bahasa daerah, dan bahasa campuran.
*   **Kekurangan:** Memerlukan kuota API, bergantung pada koneksi internet, dan pemrosesan sedikit lebih lambat dibanding Lexicon.
"""

if hasattr(st, "popover"):
    with st.sidebar.popover(":material/info: Info Metodologi (Lexicon vs LLM)", use_container_width=True):
        st.markdown(methodology_md)
else:
    with st.sidebar.expander(":material/info: Info Metodologi (Lexicon vs LLM)"):
        st.markdown(methodology_md)

if menu_selection == "Analisis Perbandingan Global":
    st.markdown("<h1><span style='color:#3498db'>SEMAN</span><span style='color:#2ecc71'>TIKA</span> : Perbandingan Global</h1>", unsafe_allow_html=True)
    if APP_MODE == "production":
        st.markdown(
            '<div style="text-align: right; margin-top: -45px; margin-bottom: 20px;">'
            '<span style="background-color: #d1fae5; color: #065f46; font-size: 0.85rem; font-weight: 700; '
            'padding: 4px 10px; border-radius: 9999px; border: 1px solid #a7f3d0;">'
            'Mode: Production (Cloud Sync)'
            '</span></div>',
            unsafe_allow_html=True
        )
    else:
        st.markdown(
            '<div style="text-align: right; margin-top: -45px; margin-bottom: 20px;">'
            '<span style="background-color: #fee2e2; color: #991b1b; font-size: 0.85rem; font-weight: 700; '
            'padding: 4px 10px; border-radius: 9999px; border: 1px solid #fecaca;">'
            'Mode: Development (Offline Lokal)'
            '</span></div>',
            unsafe_allow_html=True
        )
    st.markdown("Halaman analisis akumulatif yang menggabungkan seluruh atau sebagian riwayat video untuk perbandingan akurasi jangka panjang.")
    st.markdown("---")
    
    # 1. Bangun dictionary riwayat (gabungan Lokal + Google Sheets)
    history_videos = {}
    
    # Baca dari folder history lokal
    if os.path.exists(HISTORY_DIR):
        try:
            for f in os.listdir(HISTORY_DIR):
                if f.endswith(".csv"):
                    fpath = os.path.join(HISTORY_DIR, f)
                    df_temp = pd.read_csv(fpath)
                    clean_name = f.replace(".csv", "")
                    history_videos[clean_name] = df_temp
        except Exception:
            pass
            
    # Baca dari Google Sheets jika production
    if APP_MODE == "production" and conn is not None:
        try:
            df_gs = load_all_gsheets_data()
            if not df_gs.empty:
                grouped = df_gs.groupby(["Video ID", "Video Title"])
                for (vid_id, vid_title), group_df in grouped:
                    display_name = f"[{vid_id}] {make_safe_filename(vid_title)}"
                    local_cols = ["No", "Comment ID", "Author", "Original Comment", "Cleaned Comment", "Lexicon Sentiment", "Lexicon Score", "LLM Sentiment Global", "LLM Reason Global", "LLM Sentiment Video", "LLM Reason Video", "LLM Model", "Language", "Ground Truth", "Lexicon Time", "LLM Time Global", "LLM Time Video"]
                    df_temp = group_df.copy()
                    df_temp["No"] = range(1, len(df_temp) + 1)
                    # Filter existing columns to match local format
                    df_temp = df_temp[[col for col in local_cols if col in df_temp.columns]]
                    history_videos[display_name] = df_temp
        except Exception:
            pass
            
    history_keys = sorted(list(history_videos.keys()), reverse=True)
    
    # Inisialisasi daftar file aktif di session state
    if "active_global_files" not in st.session_state:
        st.session_state.active_global_files = list(history_keys)
        
    # Bersihkan file yang sudah dihapus dari session state
    st.session_state.active_global_files = [f for f in st.session_state.active_global_files if f in history_keys]
    
    # Tampilkan expander filter video di halaman utama
    with st.expander("⚙️ Filter Pilihan Video (Toggle Aktivasi)", expanded=True):
        st.markdown("<small>Pilih video riwayat yang ingin Anda sertakan dalam analisis dan grafik perbandingan global:</small>", unsafe_allow_html=True)
        
        # Tombol pintasan Cepat
        col_btn1, col_btn2, _ = st.columns([1.5, 1.5, 7])
        with col_btn1:
            if st.button("Pilih Semua", key="select_all_global_btn", use_container_width=True):
                st.session_state.active_global_files = list(history_keys)
                st.rerun()
        with col_btn2:
            if st.button("Kosongkan Semua", key="deselect_all_global_btn", use_container_width=True):
                st.session_state.active_global_files = []
                st.rerun()
                
        st.markdown(" ")
        
        selected_global_files = []
        if history_keys:
            cols = st.columns(3)
            for idx, h_key in enumerate(history_keys):
                clean_name = h_key
                match = re.match(r"^\[(.*?)\] (.*)$", clean_name)
                display_name = match.group(2) if match else clean_name
                
                # Potong nama jika terlalu panjang
                if len(display_name) > 35:
                    display_name = display_name[:32] + "..."
                    
                col_idx = idx % 3
                is_checked = h_key in st.session_state.active_global_files
                
                with cols[col_idx]:
                    checked = st.checkbox(f"🎥 {display_name}", value=is_checked, key=f"chk_glob_{h_key}")
                    if checked:
                        selected_global_files.append(h_key)
            
            st.session_state.active_global_files = selected_global_files
        else:
            st.info("Belum ada riwayat analisis untuk dibandingkan.")
            
    if not selected_global_files:
        st.warning("Silakan aktifkan minimal satu file riwayat pada filter di atas untuk memulai perbandingan global.")
    else:
        dfs = []
        video_accuracies = []
        
        for h_key in selected_global_files:
            try:
                df_temp = history_videos[h_key]
                df_temp = upgrade_dataframe_schema(df_temp)
                
                clean_name = h_key
                match = re.match(r"^\[(.*?)\] (.*)$", clean_name)
                vid_title = match.group(2) if match else clean_name
                
                df_temp["Video Title"] = vid_title
                dfs.append(df_temp)
                
                # Calculate accuracy for this video individually
                df_eval_temp = df_temp.dropna(subset=["Ground Truth"]).copy()
                df_eval_temp = df_eval_temp[df_eval_temp["Ground Truth"].astype(str).str.strip().str.lower().isin(["positif", "negatif", "netral"])]
                if len(df_eval_temp) > 0:
                    y_true_temp = df_eval_temp["Ground Truth"].str.strip().str.lower()
                    y_lexicon_temp = df_eval_temp["Lexicon Sentiment"].str.strip().str.lower()
                    y_llm_g_temp = df_eval_temp["LLM Sentiment Global"].str.strip().str.lower()
                    y_llm_v_temp = df_eval_temp["LLM Sentiment Video"].str.strip().str.lower()
                    
                    lex_acc_temp = accuracy_score(y_true_temp, y_lexicon_temp)
                    llm_g_acc_temp = accuracy_score(y_true_temp, y_llm_g_temp)
                    llm_v_acc_temp = accuracy_score(y_true_temp, y_llm_v_temp)
                    
                    model_temp = "meta/llama-3.1-8b-instruct"
                    if "LLM Model" in df_temp.columns and df_temp["LLM Model"].iloc[0] is not None:
                        model_temp = df_temp["LLM Model"].iloc[0]
                    model_short = model_temp.split("/")[-1] if "/" in model_temp else model_temp
                    
                    video_accuracies.append({
                        "Video": vid_title[:30] + "..." if len(vid_title) > 30 else vid_title,
                        "Lexicon Accuracy": lex_acc_temp * 100,
                        "LLM Global Accuracy": llm_g_acc_temp * 100,
                        "LLM Video Accuracy": llm_v_acc_temp * 100,
                        "LLM Model": model_short
                    })
            except Exception as e:
                st.error(f"Gagal membaca {h_key}: {e}")
                
        if dfs:
            df_global = pd.concat(dfs, ignore_index=True)
            total_comments = len(df_global)
            
            # Filter for rows with Ground Truth for evaluation
            df_global_eval = df_global.dropna(subset=["Ground Truth"]).copy()
            df_global_eval = df_global_eval[df_global_eval["Ground Truth"].astype(str).str.strip().str.lower().isin(["positif", "negatif", "netral"])]
            total_eval = len(df_global_eval)
            
            # Render KPI metrics
            col_kpi1, col_kpi2, col_kpi3 = st.columns(3)
            with col_kpi1:
                st.markdown(
                    f"""
                    <div class="metric-card">
                        <div class="metric-title">Total Komentar Terkumpul</div>
                        <div class="metric-value">{total_comments}</div>
                        <small>Dari {len(selected_global_files)} video yang dipilih</small>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            with col_kpi2:
                st.markdown(
                    f"""
                    <div class="metric-card">
                        <div class="metric-title">Ground Truth Terisi</div>
                        <div class="metric-value">{total_eval}</div>
                        <small>Persentase: {(total_eval / total_comments * 100) if total_comments > 0 else 0:.1f}% dari total</small>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            with col_kpi3:
                if total_eval > 0:
                    y_true_g = df_global_eval["Ground Truth"].str.strip().str.lower()
                    y_lex_g = df_global_eval["Lexicon Sentiment"].str.strip().str.lower()
                    y_llm_g = df_global_eval["LLM Sentiment Global"].str.strip().str.lower()
                    y_llm_v = df_global_eval["LLM Sentiment Video"].str.strip().str.lower()
                    
                    global_lex_acc = accuracy_score(y_true_g, y_lex_g) * 100
                    global_llm_g_acc = accuracy_score(y_true_g, y_llm_g) * 100
                    global_llm_v_acc = accuracy_score(y_true_g, y_llm_v) * 100
                    
                    st.markdown(
                       f"""
                       <div class="metric-card">
                           <div class="metric-title">Rata-rata Akurasi Global</div>
                           <div class="metric-value" style="font-size: 1.25rem; font-weight: 700; color: #1e293b; line-height: 1.4;">
                               Lexicon: {global_lex_acc:.1f}%<br/>
                               LLM Global: {global_llm_g_acc:.1f}%<br/>
                               LLM Video: {global_llm_v_acc:.1f}%
                           </div>
                       </div>
                       """,
                       unsafe_allow_html=True
                    )
                else:
                    st.markdown(
                        """
                        <div class="metric-card">
                            <div class="metric-title">Rata-rata Akurasi Global</div>
                            <div class="metric-value" style="font-size: 1.3rem; color: #64748b;">N/A</div>
                            <small>Isi Ground Truth terlebih dahulu</small>
                        </div>
                        """,
                        unsafe_allow_html=True
                    )
                    
            if total_eval == 0:
                st.warning("Belum ada data Ground Truth yang diisi di seluruh video yang dipilih. Metrik komparasi tidak dapat ditampilkan.")
            else:
                # Plot global donut charts
                st.subheader("📊 Sebaran Sentimen Akumulatif")
                
                y_true_g = df_global_eval["Ground Truth"].str.strip().str.lower()
                y_lex_g = df_global_eval["Lexicon Sentiment"].str.strip().str.lower()
                y_llm_g = df_global_eval["LLM Sentiment Global"].str.strip().str.lower()
                y_llm_v = df_global_eval["LLM Sentiment Video"].str.strip().str.lower()
                
                sentiment_labels = ["positif", "negatif", "netral"]
                color_map = {"positif": "#2ecc71", "negatif": "#e74c3c", "netral": "#95a5a6"}
                
                def get_sizes_and_colors(series):
                    counts = series.value_counts()
                    sizes = []
                    colors = []
                    labels = []
                    for label in sentiment_labels:
                        count = counts.get(label, 0)
                        if count > 0:
                            sizes.append(count)
                            colors.append(color_map[label])
                            labels.append(label.capitalize())
                    return sizes, colors, labels
                    
                gt_sizes, gt_colors, gt_labels = get_sizes_and_colors(y_true_g)
                lex_sizes, lex_colors, lex_labels = get_sizes_and_colors(y_lex_g)
                llm_g_sizes, llm_g_colors, llm_g_labels = get_sizes_and_colors(y_llm_g)
                llm_v_sizes, llm_v_colors, llm_v_labels = get_sizes_and_colors(y_llm_v)
                
                fig_donut_g, axs = plt.subplots(1, 4, figsize=(22, 6))
                
                # Donut 1: GT
                if gt_sizes:
                    axs[0].pie(gt_sizes, labels=gt_labels, autopct='%1.1f%%', startangle=90, colors=gt_colors, pctdistance=0.75, textprops=dict(color="black", weight="bold"))
                    axs[0].add_artist(plt.Circle((0,0), 0.50, fc='white'))
                    axs[0].set_title("Global Ground Truth", fontsize=12, weight="bold")
                else:
                    axs[0].text(0.5, 0.5, 'Tidak ada data', ha='center', va='center')
                    
                # Donut 2: Lexicon
                if lex_sizes:
                    axs[1].pie(lex_sizes, labels=lex_labels, autopct='%1.1f%%', startangle=90, colors=lex_colors, pctdistance=0.75, textprops=dict(color="black", weight="bold"))
                    axs[1].add_artist(plt.Circle((0,0), 0.50, fc='white'))
                    axs[1].set_title("Global Lexicon-based", fontsize=12, weight="bold")
                else:
                    axs[1].text(0.5, 0.5, 'Tidak ada data', ha='center', va='center')
                    
                # Donut 3: LLM Global
                if llm_g_sizes:
                    axs[2].pie(llm_g_sizes, labels=llm_g_labels, autopct='%1.1f%%', startangle=90, colors=llm_g_colors, pctdistance=0.75, textprops=dict(color="black", weight="bold"))
                    axs[2].add_artist(plt.Circle((0,0), 0.50, fc='white'))
                    axs[2].set_title("Global LLM Konteks Global", fontsize=12, weight="bold")
                else:
                    axs[2].text(0.5, 0.5, 'Tidak ada data', ha='center', va='center')
                    
                # Donut 4: LLM Video
                if llm_v_sizes:
                    axs[3].pie(llm_v_sizes, labels=llm_v_labels, autopct='%1.1f%%', startangle=90, colors=llm_v_colors, pctdistance=0.75, textprops=dict(color="black", weight="bold"))
                    axs[3].add_artist(plt.Circle((0,0), 0.50, fc='white'))
                    axs[3].set_title("Global LLM Konteks Video", fontsize=12, weight="bold")
                else:
                    axs[3].text(0.5, 0.5, 'Tidak ada data', ha='center', va='center')
                    
                plt.tight_layout()
                st.pyplot(fig_donut_g)
                plt.close()
                
                # Tabs
                tab_glob1, tab_glob2 = st.tabs(["📊 Global Metrics Comparison", "📈 Per Video Accuracy Comparison"])
                
                with tab_glob1:
                    st.markdown("### Perbandingan Metrik Evaluasi Akumulatif (Global)")
                    
                    global_lex_acc = accuracy_score(y_true_g, y_lex_g)
                    global_lex_prec, global_lex_rec, global_lex_f1, _ = precision_recall_fscore_support(y_true_g, y_lex_g, average='macro', zero_division=0)
                    
                    global_llm_g_acc = accuracy_score(y_true_g, y_llm_g)
                    global_llm_g_prec, global_llm_g_rec, global_llm_g_f1, _ = precision_recall_fscore_support(y_true_g, y_llm_g, average='macro', zero_division=0)
                    
                    global_llm_v_acc = accuracy_score(y_true_g, y_llm_v)
                    global_llm_v_prec, global_llm_v_rec, global_llm_v_f1, _ = precision_recall_fscore_support(y_true_g, y_llm_v, average='macro', zero_division=0)
                    
                    metrics_g = ['Akurasi', 'Presisi', 'Recall', 'F1-Score']
                    lex_scores_g = [global_lex_acc * 100, global_lex_prec * 100, global_lex_rec * 100, global_lex_f1 * 100]
                    llm_g_scores_g = [global_llm_g_acc * 100, global_llm_g_prec * 100, global_llm_g_rec * 100, global_llm_g_f1 * 100]
                    llm_v_scores_g = [global_llm_v_acc * 100, global_llm_v_prec * 100, global_llm_v_rec * 100, global_llm_v_f1 * 100]
                    
                    x_g = np.arange(len(metrics_g))
                    width_g = 0.25
                    
                    fig_metrics_g, ax_mg = plt.subplots(figsize=(12, 5))
                    rects_g1 = ax_mg.bar(x_g - width_g, lex_scores_g, width_g, label='Lexicon-Based', color='#3498db')
                    rects_g2 = ax_mg.bar(x_g, llm_g_scores_g, width_g, label='LLM Konteks Global', color='#e67e22')
                    rects_g3 = ax_mg.bar(x_g + width_g, llm_v_scores_g, width_g, label='LLM Konteks ke Video', color='#2ecc71')
                    
                    ax_mg.set_ylabel('Skor (%)', weight="bold")
                    ax_mg.set_title('Perbandingan Metrik Evaluasi Akumulatif (Global)', weight="bold", fontsize=12)
                    ax_mg.set_xticks(x_g)
                    ax_mg.set_xticklabels(metrics_g, weight="bold")
                    ax_mg.set_ylim(0, 115)
                    ax_mg.legend()
                    
                    def autolabel_g(rects):
                        for rect in rects:
                            height = rect.get_height()
                            ax_mg.annotate(f'{height:.1f}%',
                                        xy=(rect.get_x() + rect.get_width() / 2, height),
                                        xytext=(0, 3),
                                        textcoords="offset points",
                                        ha='center', va='bottom', weight="bold", fontsize=8)
                    autolabel_g(rects_g1)
                    autolabel_g(rects_g2)
                    autolabel_g(rects_g3)
                    
                    plt.tight_layout()
                    st.pyplot(fig_metrics_g)
                    plt.close()
                    
                with tab_glob2:
                    if video_accuracies:
                        st.markdown("### Perbandingan Akurasi antara Lexicon dan LLM untuk Setiap Video")
                        df_vid_acc = pd.DataFrame(video_accuracies)
                        
                        fig_line, ax_l = plt.subplots(figsize=(12, 5))
                        x_indices = np.arange(len(df_vid_acc))
                        
                        ax_l.plot(x_indices, df_vid_acc["Lexicon Accuracy"], marker='o', linewidth=2, color='#3498db', label='Lexicon Accuracy')
                        ax_l.plot(x_indices, df_vid_acc["LLM Global Accuracy"], marker='^', linewidth=2, color='#e67e22', label='LLM Global Accuracy')
                        ax_l.plot(x_indices, df_vid_acc["LLM Video Accuracy"], marker='s', linewidth=2, color='#2ecc71', label='LLM Video Accuracy')
                        
                        for i in range(len(df_vid_acc)):
                            ax_l.annotate(f'{df_vid_acc["Lexicon Accuracy"].iloc[i]:.1f}%', (x_indices[i], df_vid_acc["Lexicon Accuracy"].iloc[i]), textcoords="offset points", xytext=(0,10), ha='center', color='#1e3a8a', weight="bold", fontsize=8)
                            ax_l.annotate(f'{df_vid_acc["LLM Global Accuracy"].iloc[i]:.1f}%', (x_indices[i], df_vid_acc["LLM Global Accuracy"].iloc[i]), textcoords="offset points", xytext=(0,10), ha='center', color='#d35400', weight="bold", fontsize=8)
                            ax_l.annotate(f'{df_vid_acc["LLM Video Accuracy"].iloc[i]:.1f}%', (x_indices[i], df_vid_acc["LLM Video Accuracy"].iloc[i]), textcoords="offset points", xytext=(0,-15), ha='center', color='#166534', weight="bold", fontsize=8)
                            
                        ax_l.set_xticks(x_indices)
                        ax_l.set_xticklabels(df_vid_acc["Video"], rotation=30, ha='right', weight="bold", fontsize=9)
                        ax_l.set_ylabel('Akurasi (%)', weight="bold")
                        ax_l.set_ylim(0, 115)
                        ax_l.set_title('Akurasi Metode Lexicon vs LLM per Video', weight="bold", fontsize=12)
                        ax_l.grid(True, linestyle='--', alpha=0.5)
                        ax_l.legend()
                        
                        plt.tight_layout()
                        st.pyplot(fig_line)
                        plt.close()
                        
                        # Display table
                        st.dataframe(
                            df_vid_acc,
                            column_config={
                                "Video": st.column_config.TextColumn("Judul Video", width="large"),
                                "Lexicon Accuracy": st.column_config.NumberColumn("Akurasi Lexicon", format="%.2f%%"),
                                "LLM Global Accuracy": st.column_config.NumberColumn("Akurasi LLM Global", format="%.2f%%"),
                                "LLM Video Accuracy": st.column_config.NumberColumn("Akurasi LLM Video", format="%.2f%%"),
                                "LLM Model": st.column_config.TextColumn("Model LLM yang Digunakan")
                            },
                            use_container_width=True,
                            hide_index=True
                        )
                    else:
                        st.info("Belum ada video dengan Ground Truth terisi untuk dibandingkan.")
    st.stop()
# Main Dashboard Area (SEMANTIKA)
st.markdown("<h1><span style='color:#3498db'>SEMAN</span><span style='color:#2ecc71'>TIKA</span> : Sentiment Analysis Dashboard</h1>", unsafe_allow_html=True)
if APP_MODE == "production":
    st.markdown(
        '<div style="text-align: right; margin-top: -45px; margin-bottom: 20px;">'
        '<span style="background-color: #d1fae5; color: #065f46; font-size: 0.85rem; font-weight: 700; '
        'padding: 4px 10px; border-radius: 9999px; border: 1px solid #a7f3d0;">'
        'Mode: Production (Cloud Sync)'
        '</span></div>',
        unsafe_allow_html=True
    )
else:
    st.markdown(
        '<div style="text-align: right; margin-top: -45px; margin-bottom: 20px;">'
        '<span style="background-color: #fee2e2; color: #991b1b; font-size: 0.85rem; font-weight: 700; '
        'padding: 4px 10px; border-radius: 9999px; border: 1px solid #fecaca;">'
        'Mode: Development (Offline Lokal)'
        '</span></div>',
        unsafe_allow_html=True
    )
st.markdown("Aplikasi perbandingan performa analisis sentimen berbasis **Lexicon-based (Sastrawi + InSet)** dan **LLM-based (NVIDIA NIM Llama 3.1)**.")
st.markdown("---")

if st.session_state.df is not None:
    # Header: Video Info
    st.markdown(f"### :material/movie: **{st.session_state.video_title}**")
    col_hdr1, col_hdr2, col_hdr3, col_hdr4 = st.columns(4)
    with col_hdr1:
        st.markdown(f":material/link: **Link Video:** [{st.session_state.video_url}]({st.session_state.video_url})")
    with col_hdr2:
        st.markdown(f":material/neurology: **Model LLM Aktif:** `{st.session_state.llm_model}`")
    with col_hdr3:
        if "Language" in st.session_state.df.columns:
            lang_counts = st.session_state.df["Language"].value_counts()
            total_comments = len(st.session_state.df)
            lang_labels = []
            for lang, count in lang_counts.items():
                pct = (count / total_comments * 100) if total_comments > 0 else 0
                if str(lang).strip().lower() == "id":
                    lang_labels.append(f"Indonesia (ID) {pct:.1f}%")
                elif str(lang).strip().lower() == "en":
                    lang_labels.append(f"Inggris (EN) {pct:.1f}%")
                else:
                    lang_labels.append(f"{str(lang).upper()} {pct:.1f}%")
            lang_label = " & ".join(lang_labels) if lang_labels else "Indonesia (ID)"
        else:
            lang_label = "Inggris (EN)" if st.session_state.detected_lang == "en" else "Indonesia (ID)"
        st.markdown(f":material/translate: **Bahasa Terdeteksi:** `{lang_label}`")
    with col_hdr4:
        st.markdown(f":material/settings: **Mode Analisis:** `Dual Mode (Global & Video)`")
    
    st.markdown("---")
    
    # Root Tabs
    tab_global, tab_video, tab_compare = st.tabs([
        "🌐 Mode Konteks Global",
        "🎥 Mode Konteks ke Video",
        "📊 Perbandingan Performa"
    ])
    
    with tab_global:
        render_mode_tab_content("global", "LLM Sentiment Global", "LLM Reason Global", "Konteks Global")
        
    with tab_video:
        render_mode_tab_content("video", "LLM Sentiment Video", "LLM Reason Video", "Konteks ke Video")
        
    with tab_compare:
        st.markdown("### Perbandingan Performa: Konteks Global vs Konteks ke Video")
        
        df_eval = st.session_state.df.dropna(subset=["Ground Truth"]).copy()
        df_eval = df_eval[df_eval["Ground Truth"].astype(str).str.strip().str.lower().isin(["positif", "negatif", "netral"])]
        
        if len(df_eval) == 0:
            st.warning("Belum ada Ground Truth yang diisi. Silakan isi beberapa baris pada kolom Ground Truth di salah satu tab mode untuk menampilkan perbandingan performa.", icon=":material/warning:")
        else:
            y_true = df_eval["Ground Truth"].str.strip().str.lower()
            
            # Scores for Lexicon
            y_lexicon = df_eval["Lexicon Sentiment"].str.strip().str.lower()
            lex_acc = accuracy_score(y_true, y_lexicon)
            lex_prec, lex_rec, lex_f1, _ = precision_recall_fscore_support(y_true, y_lexicon, average='macro', zero_division=0)
            
            # Scores for LLM Global
            y_llm_global = df_eval["LLM Sentiment Global"].str.strip().str.lower()
            llm_g_acc = accuracy_score(y_true, y_llm_global)
            llm_g_prec, llm_g_rec, llm_g_f1, _ = precision_recall_fscore_support(y_true, y_llm_global, average='macro', zero_division=0)
            
            # Scores for LLM Video
            y_llm_video = df_eval["LLM Sentiment Video"].str.strip().str.lower()
            llm_v_acc = accuracy_score(y_true, y_llm_video)
            llm_v_prec, llm_v_rec, llm_v_f1, _ = precision_recall_fscore_support(y_true, y_llm_video, average='macro', zero_division=0)
            
            # 1. Bar Chart Comparison
            st.markdown("#### Perbandingan Akurasi & F1-Score")
            metrics = ["Akurasi", "F1-Score"]
            lex_scores = [lex_acc * 100, lex_f1 * 100]
            llm_g_scores = [llm_g_acc * 100, llm_g_f1 * 100]
            llm_v_scores = [llm_v_acc * 100, llm_v_f1 * 100]
            
            x = np.arange(len(metrics))
            width = 0.25
            
            fig_cmp, ax_cmp = plt.subplots(figsize=(10, 5))
            rects1 = ax_cmp.bar(x - width, lex_scores, width, label='Lexicon-Based', color='#3498db')
            rects2 = ax_cmp.bar(x, llm_g_scores, width, label='LLM Konteks Global', color='#e67e22')
            rects3 = ax_cmp.bar(x + width, llm_v_scores, width, label='LLM Konteks ke Video', color='#2ecc71')
            
            ax_cmp.set_ylabel('Skor (%)', weight="bold")
            ax_cmp.set_title('Perbandingan Metrik: Lexicon vs LLM Global vs LLM Video', weight="bold", fontsize=12)
            ax_cmp.set_xticks(x)
            ax_cmp.set_xticklabels(metrics, weight="bold")
            ax_cmp.set_ylim(0, 115)
            ax_cmp.legend()
            
            def autolabel_cmp(rects):
                for rect in rects:
                    height = rect.get_height()
                    ax_cmp.annotate(f'{height:.1f}%',
                                xy=(rect.get_x() + rect.get_width() / 2, height),
                                xytext=(0, 3),
                                textcoords="offset points",
                                ha='center', va='bottom', weight="bold", fontsize=8)
            autolabel_cmp(rects1)
            autolabel_cmp(rects2)
            autolabel_cmp(rects3)
            
            plt.tight_layout()
            st.pyplot(fig_cmp)
            plt.close()
            
            # 2. Side-by-side Point Cards
            st.markdown("#### Perbandingan Skor Poin")
            lex_points, llm_g_points, llm_v_points = 0, 0, 0
            for idx, row in df_eval.iterrows():
                gt = str(row["Ground Truth"]).strip().lower()
                lex = str(row["Lexicon Sentiment"]).strip().lower()
                llm_g = str(row["LLM Sentiment Global"]).strip().lower()
                llm_v = str(row["LLM Sentiment Video"]).strip().lower()
                
                lex_points += 1 if lex == gt else -1
                llm_g_points += 1 if llm_g == gt else -1
                llm_v_points += 1 if llm_v == gt else -1
                
            col_c1, col_c2, col_c3 = st.columns(3)
            with col_c1:
                st.markdown(
                    f"""
                    <div class="point-card lexicon-card">
                        <h3>LEXICON-BASED</h3>
                        <div style="font-size: 2.5rem; font-weight: 800; margin: 5px 0;">{lex_points} Poin</div>
                        <p style="margin-top: 5px; font-size: 0.8rem;">Sastrawi + InSet</p>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            with col_c2:
                st.markdown(
                    f"""
                    <div class="point-card llm-card" style="background: linear-gradient(135deg, #d35400, #e67e22);">
                        <h3>LLM KONTEKS GLOBAL</h3>
                        <div style="font-size: 2.5rem; font-weight: 800; margin: 5px 0;">{llm_g_points} Poin</div>
                        <p style="margin-top: 5px; font-size: 0.8rem;">Llama 3.1 8B (NIM)</p>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            with col_c3:
                st.markdown(
                    f"""
                    <div class="point-card llm-card">
                        <h3>LLM KONTEKS VIDEO</h3>
                        <div style="font-size: 2.5rem; font-weight: 800; margin: 5px 0;">{llm_v_points} Poin</div>
                        <p style="margin-top: 5px; font-size: 0.8rem;">Llama 3.1 8B (NIM)</p>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
                
            # 3. Summary Table
            st.markdown("#### Tabel Perbandingan Seluruh Metrik")
            summary_data = {
                "Metrik": ["Akurasi", "Presisi (Macro)", "Recall (Macro)", "F1-Score (Macro)"],
                "Lexicon-Based": [f"{lex_acc*100:.1f}%", f"{lex_prec*100:.1f}%", f"{lex_rec*100:.1f}%", f"{lex_f1*100:.1f}%"],
                "LLM Konteks Global": [f"{llm_g_acc*100:.1f}%", f"{llm_g_prec*100:.1f}%", f"{llm_g_rec*100:.1f}%", f"{llm_g_f1*100:.1f}%"],
                "LLM Konteks ke Video": [f"{llm_v_acc*100:.1f}%", f"{llm_v_prec*100:.1f}%", f"{llm_v_rec*100:.1f}%", f"{llm_v_f1*100:.1f}%"]
            }
            st.table(pd.DataFrame(summary_data))

elif menu_selection == "Kelola Kamus Slang":
    st.markdown("<h1><span style='color:#3498db'>SEMAN</span><span style='color:#2ecc71'>TIKA</span> : Kelola Kamus Slang</h1>", unsafe_allow_html=True)
    st.markdown("Halaman ini digunakan untuk mengelola kamus singkatan (*slang*) yang digunakan dalam pra-pemrosesan metode Lexicon.")
    st.markdown("---")
    
    from src.normalizer import load_slang_dict, save_custom_slang, CUSTOM_SLANG_PATH
    
    # We want to edit the custom slang dictionary.
    import json
    custom_slang = {}
    if os.path.exists(CUSTOM_SLANG_PATH):
        try:
            with open(CUSTOM_SLANG_PATH, "r", encoding="utf-8") as f:
                custom_slang = json.load(f)
        except Exception:
            pass
            
    # Convert custom slang dict to dataframe for editing
    slang_data = [{"Singkatan (Slang)": k, "Kata Baku": v} for k, v in custom_slang.items()]
    df_slang = pd.DataFrame(slang_data)
    if df_slang.empty:
        df_slang = pd.DataFrame(columns=["Singkatan (Slang)", "Kata Baku"])
        
    st.info("Kamus slang default (seperti 'yg' -> 'yang', 'bgt' -> 'banget') sudah aktif di sistem secara otomatis. "
            "Gunakan tabel di bawah ini untuk **menambahkan singkatan baru** atau **meng-override** kata slang default.", icon=":material/info:")
            
    edited_slang_df = st.data_editor(
        df_slang,
        num_rows="dynamic",
        use_container_width=True,
        column_config={
            "Singkatan (Slang)": st.column_config.TextColumn("Singkatan (Slang) / Kata Gaul", help="Kata gaul/singkatan yang ingin dicocokkan (contoh: 'mager')", required=True),
            "Kata Baku": st.column_config.TextColumn("Kata Baku / Normalisasi", help="Kata baku hasil konversinya (contoh: 'malas')", required=True),
        },
        key="slang_editor"
    )
    
    col1, col2 = st.columns([1, 4])
    with col1:
        if st.button(":material/save: Simpan Perubahan", use_container_width=True):
            # Convert dataframe back to dictionary
            new_custom_slang = {}
            for _, row in edited_slang_df.iterrows():
                k = str(row.get("Singkatan (Slang)", "")).strip().lower()
                v = str(row.get("Kata Baku", "")).strip().lower()
                if k:
                    new_custom_slang[k] = v
            if save_custom_slang(new_custom_slang):
                st.success("Kamus slang kustom berhasil disimpan dan diperbarui!")
                st.rerun()
            else:
                st.error("Gagal menyimpan kamus slang.")

else:
    # Welcome message with clean CSS styling
    st.markdown(
        """
        <div class="info-box" style="border-left-color: #3498db;">
            <h3>Selamat datang di SEMANTIKA!</h3>
            <p>Silakan gunakan panel konfigurasi di sidebar kiri untuk menghubungkan dashboard dengan YouTube dan memulai analisis sentimen.</p>
        </div>
        """,
        unsafe_allow_html=True
    )
    st.markdown("""
    ### Langkah Memulai:
    1. Pastikan file `.env` Anda sudah terisi dengan **NVIDIA API Key** yang valid.
    2. Masukkan URL video YouTube/Shorts di panel konfigurasi sebelah kiri.
    3. Tentukan batas jumlah komentar (contoh: 100).
    4. Pilih model LLM yang ingin digunakan (disarankan: `meta/llama-3.1-8b-instruct` untuk pemrosesan cepat).
    5. Klik tombol **Mulai Analisis Data**.
    
    Aplikasi akan mengunduh komentar, menganalisis dengan kedua metode, dan menampilkan tabel interaktif untuk pengisian **Ground Truth**.
    """)
